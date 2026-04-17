from __future__ import annotations

from pathlib import Path
from datetime import datetime
import re
import hashlib
import subprocess
import shutil
from urllib.parse import quote

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt, RGBColor as DocRGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt as PptPt


class DocumentExportService:
    """Builds PPTX and DOCX files from structured lesson content."""

    TARGET_SLIDE_WIDTH = Inches(13.333)
    TARGET_SLIDE_HEIGHT = Inches(7.5)
    TEMPLATE_MIN_BODY_LINES = 2
    TEMPLATE_MIN_CHAR_CAPACITY = 80
    PLACEHOLDER_HINTS = (
        "add a main point",
        "briefly elaborate",
        "elaborate on",
        "lorem ipsum",
        "click to add",
        "title",
        "subtitle",
    )

    THEMES = {
        "科技风": {"bg": (15, 24, 49), "soft": (26, 39, 74), "card": (244, 250, 255), "accent": (56, 194, 255), "text": (244, 248, 255), "dark_text": (23, 35, 57), "family": "neon"},
        "极简风": {"bg": (247, 245, 240), "soft": (255, 252, 247), "card": (255, 255, 255), "accent": (191, 112, 64), "text": (35, 31, 28), "dark_text": (35, 31, 28), "family": "minimal"},
        "学术风": {"bg": (248, 250, 253), "soft": (238, 244, 252), "card": (255, 255, 255), "accent": (30, 83, 160), "text": (26, 33, 44), "dark_text": (26, 33, 44), "family": "formal"},
        "卡通风": {"bg": (255, 249, 235), "soft": (255, 255, 246), "card": (255, 255, 255), "accent": (245, 153, 49), "text": (73, 54, 24), "dark_text": (73, 54, 24), "family": "playful"},
        "清新风": {"bg": (242, 250, 245), "soft": (250, 255, 251), "card": (255, 255, 255), "accent": (56, 152, 113), "text": (33, 57, 43), "dark_text": (33, 57, 43), "family": "airy"},
        "中国风": {"bg": (247, 243, 236), "soft": (253, 250, 246), "card": (255, 252, 248), "accent": (155, 48, 48), "text": (72, 46, 43), "dark_text": (72, 46, 43), "family": "elegant"},
        "商务风": {"bg": (241, 245, 250), "soft": (248, 251, 255), "card": (255, 255, 255), "accent": (37, 78, 138), "text": (28, 39, 58), "dark_text": (28, 39, 58), "family": "split"},
        "可爱风": {"bg": (255, 244, 248), "soft": (255, 250, 252), "card": (255, 255, 255), "accent": (232, 103, 153), "text": (87, 45, 64), "dark_text": (87, 45, 64), "family": "playful"},
        "莫兰迪风": {"bg": (237, 233, 227), "soft": (247, 244, 240), "card": (251, 249, 246), "accent": (136, 124, 114), "text": (72, 66, 61), "dark_text": (72, 66, 61), "family": "minimal"},
        "新中式": {"bg": (244, 238, 230), "soft": (250, 245, 239), "card": (254, 251, 247), "accent": (129, 55, 45), "text": (72, 45, 40), "dark_text": (72, 45, 40), "family": "elegant"},
    "暗夜风": {"bg": (244, 246, 252), "soft": (250, 251, 255), "card": (255, 255, 255), "accent": (99, 102, 241), "text": (31, 41, 55), "dark_text": (31, 41, 55), "family": "formal"},
    "霓虹风": {"bg": (245, 247, 255), "soft": (251, 252, 255), "card": (255, 255, 255), "accent": (236, 72, 153), "text": (55, 48, 163), "dark_text": (55, 48, 163), "family": "gradient"},
        "杂志风": {"bg": (250, 247, 240), "soft": (255, 253, 248), "card": (255, 255, 255), "accent": (34, 34, 34), "text": (28, 28, 28), "dark_text": (28, 28, 28), "family": "editorial"},
    "黑金风": {"bg": (250, 246, 236), "soft": (255, 252, 244), "card": (255, 255, 255), "accent": (180, 120, 32), "text": (72, 52, 20), "dark_text": (72, 52, 20), "family": "formal"},
        "粉彩风": {"bg": (252, 244, 247), "soft": (255, 251, 253), "card": (255, 255, 255), "accent": (241, 164, 186), "text": (92, 71, 78), "dark_text": (92, 71, 78), "family": "airy"},
        "自然风": {"bg": (244, 247, 238), "soft": (250, 252, 246), "card": (255, 255, 255), "accent": (126, 153, 80), "text": (53, 69, 41), "dark_text": (53, 69, 41), "family": "airy"},
        "森林风": {"bg": (231, 241, 233), "soft": (245, 250, 245), "card": (252, 255, 252), "accent": (59, 110, 75), "text": (35, 58, 43), "dark_text": (35, 58, 43), "family": "split"},
        "海洋风": {"bg": (231, 244, 251), "soft": (246, 251, 255), "card": (255, 255, 255), "accent": (39, 124, 172), "text": (30, 65, 90), "dark_text": (30, 65, 90), "family": "airy"},
        "复古风": {"bg": (244, 232, 214), "soft": (251, 245, 235), "card": (255, 250, 242), "accent": (153, 92, 53), "text": (81, 57, 40), "dark_text": (81, 57, 40), "family": "editorial"},
        "手账风": {"bg": (251, 247, 238), "soft": (255, 252, 247), "card": (255, 255, 252), "accent": (217, 130, 95), "text": (93, 70, 58), "dark_text": (93, 70, 58), "family": "playful"},
        "未来感": {"bg": (11, 23, 43), "soft": (26, 41, 70), "card": (36, 53, 88), "accent": (0, 220, 255), "text": (238, 249, 255), "dark_text": (238, 249, 255), "family": "neon"},
        "插画风": {"bg": (255, 249, 241), "soft": (255, 253, 248), "card": (255, 255, 255), "accent": (238, 141, 69), "text": (79, 60, 44), "dark_text": (79, 60, 44), "family": "playful"},
        "教育蓝": {"bg": (240, 246, 255), "soft": (248, 251, 255), "card": (255, 255, 255), "accent": (37, 99, 235), "text": (30, 58, 95), "dark_text": (30, 58, 95), "family": "formal"},
        "暖阳风": {"bg": (255, 248, 236), "soft": (255, 252, 246), "card": (255, 255, 255), "accent": (236, 153, 31), "text": (88, 61, 30), "dark_text": (88, 61, 30), "family": "airy"},
        "几何风": {"bg": (245, 247, 250), "soft": (251, 252, 255), "card": (255, 255, 255), "accent": (112, 87, 255), "text": (43, 47, 62), "dark_text": (43, 47, 62), "family": "geometry"},
        "玻璃拟态": {"bg": (231, 241, 247), "soft": (243, 248, 252), "card": (255, 255, 255), "accent": (89, 138, 199), "text": (35, 52, 71), "dark_text": (35, 52, 71), "family": "glass"},
        "扁平风": {"bg": (245, 247, 250), "soft": (250, 252, 255), "card": (255, 255, 255), "accent": (39, 112, 255), "text": (34, 43, 58), "dark_text": (34, 43, 58), "family": "flat"},
        "渐变风": {"bg": (242, 238, 255), "soft": (250, 247, 255), "card": (255, 255, 255), "accent": (140, 92, 246), "text": (52, 38, 84), "dark_text": (52, 38, 84), "family": "gradient"},
    }

    def __init__(self, output_dir: Path):
        self.output_dir = output_dir
        self.theme = self.THEMES["教育蓝"]
        self.template_dir = output_dir.parent / "ppt_format"
        self.preview_dir = output_dir.parent / "static" / "generated_ppt_previews"
        self.preview_dir.mkdir(parents=True, exist_ok=True)
        self.external_templates = self._discover_external_templates()

    def _discover_external_templates(self):
        if not self.template_dir.exists():
            return {}
        templates = {}
        for file in sorted(self.template_dir.glob("*.pptx")):
            templates[file.stem.strip()] = file
        return templates

    def build_package(self, package: dict):
        ppt_name = self._build_export_filename(package, kind="课件", ext="pptx")
        doc_name = self._build_export_filename(package, kind="教案", ext="docx")
        self._build_ppt(self.output_dir / ppt_name, package)
        self._build_doc(self.output_dir / doc_name, package)
        return {"pptx": ppt_name, "docx": doc_name}

    def build_ppt(self, package: dict):
        ppt_name = self._build_export_filename(package, kind="课件", ext="pptx")
        self._build_ppt(self.output_dir / ppt_name, package)
        return {"pptx": ppt_name}

    def build_doc(self, package: dict):
        doc_name = self._build_export_filename(package, kind="教案", ext="docx")
        self._build_doc(self.output_dir / doc_name, package)
        return {"docx": doc_name}

    def _build_export_filename(self, package: dict, kind: str, ext: str) -> str:
        summary = package.get("summary") or {}
        raw_topic = str(summary.get("course_theme") or "").strip()
        raw_title = str(package.get("title") or "").strip()
        raw_knowledge = str(summary.get("knowledge_points") or "").strip()
        base_name = raw_topic or raw_title or raw_knowledge or "教学内容"
        safe_base = self._sanitize_filename(base_name)
        timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
        stem = f"{safe_base}-{kind}-{timestamp}"
        candidate = f"{stem}.{ext}"
        index = 2
        while (self.output_dir / candidate).exists():
            candidate = f"{stem}-{index}.{ext}"
            index += 1
        return candidate

    def _sanitize_filename(self, value: str, max_length: int = 36) -> str:
        cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]', " ", str(value or "").strip())
        cleaned = re.sub(r"\s+", " ", cleaned).strip(" .")
        if not cleaned:
            return "教学内容"
        return cleaned[:max_length].rstrip(" .") or "教学内容"

    def _build_ppt(self, file_path: Path, package: dict):
        style_signal = package.get("selected_template") or package.get("theme_style") or package.get("summary", {}).get("style") or "教育蓝"
        self._apply_theme(style_signal)
        slides = self._normalize_slides_for_ppt(package.get("slides", []))
        template_name = str(package.get("selected_template") or "").strip()
        template_path = self.external_templates.get(template_name)

        if template_path and template_path.exists():
            prs = Presentation(str(template_path))
            prs.slide_width = self.TARGET_SLIDE_WIDTH
            prs.slide_height = self.TARGET_SLIDE_HEIGHT
            template_applied = self._apply_content_to_template(prs, package, slides)
            if template_applied:
                prs.save(str(file_path))
                return

        prs = Presentation()
        prs.slide_width = self.TARGET_SLIDE_WIDTH
        prs.slide_height = self.TARGET_SLIDE_HEIGHT
        for slide_data in slides:
            layout = slide_data.get("layout", "content")
            if layout == "cover":
                self._add_cover_slide(prs, package)
            elif layout == "agenda":
                self._add_agenda_slide(prs, slide_data)
            elif layout == "case":
                self._add_case_slide(prs, slide_data)
            elif layout == "summary":
                self._add_summary_slide(prs, slide_data, package)
            elif layout == "interactive":
                self._add_interactive_slide(prs, slide_data)
            else:
                self._add_content_slide(prs, slide_data)
        prs.save(str(file_path))

    def _apply_content_to_template(self, prs: Presentation, package: dict, slides: list[dict]) -> bool:
        if not slides:
            return True

        if len(prs.slide_layouts) == 0:
            return False

        if len(prs.slides) == 0:
            default_layout = prs.slide_layouts[0]
            for _ in slides:
                prs.slides.add_slide(default_layout)
        elif len(prs.slides) < len(slides):
            for index in range(len(prs.slides), len(slides)):
                layout = prs.slide_layouts[min(index, len(prs.slide_layouts) - 1)]
                prs.slides.add_slide(layout)

        for index, slide_data in enumerate(slides):
            text_lines = self._build_template_text_lines(slide_data, package, index=index, total=len(slides))
            if not self._template_slide_has_enough_capacity(prs.slides[index], text_lines):
                return False
            self._fill_template_slide_text(prs.slides[index], text_lines)

        for index in range(len(prs.slides) - 1, len(slides) - 1, -1):
            self._remove_slide_at(prs, index)

        return True

    def _template_slide_has_enough_capacity(self, slide, text_lines: list[str]) -> bool:
        title_shape, body_shapes = self._collect_text_shapes(slide)
        if not title_shape and not body_shapes:
            return False

        body_lines = [line for line in (text_lines[1:] if len(text_lines) > 1 else []) if str(line or "").strip()]
        if not body_lines:
            return True

        if not body_shapes:
            return False

        line_capacity = sum(self._estimate_shape_line_capacity(shape) for shape in body_shapes)
        char_capacity = sum(self._estimate_shape_char_capacity(shape) * self._estimate_shape_line_capacity(shape) for shape in body_shapes)
        needed_lines = min(len(body_lines), 6)
        needed_chars = sum(len(str(line)) for line in body_lines)

        if line_capacity < max(self.TEMPLATE_MIN_BODY_LINES, needed_lines):
            return False
        if char_capacity < max(self.TEMPLATE_MIN_CHAR_CAPACITY, int(needed_chars * 0.65)):
            return False
        return True

    def _build_template_text_lines(self, slide_data: dict, package: dict, *, index: int, total: int) -> list[str]:
        layout = str(slide_data.get("layout") or "content")
        title = self._compact_text(slide_data.get("title") or f"第 {index + 1} 页", max_len=40)
        bullets = [self._compact_text(item, max_len=68) for item in (slide_data.get("bullets") or [])]
        bullets = [item for item in bullets if item]
        summary = package.get("summary") or {}
        supplemental = [
            self._compact_text(f"知识点：{summary.get('knowledge_points', '')}", max_len=68),
            self._compact_text(f"重难点：{summary.get('key_difficulties', '')}", max_len=68),
        ]
        supplemental = [item for item in supplemental if item and not item.endswith("：")]

        if layout == "cover":
            return [
                package.get("title") or title,
                f"课程主题：{summary.get('course_theme', '')}",
                f"知识点：{summary.get('knowledge_points', '')}",
                f"课时安排：{summary.get('lesson_periods', '')}",
            ]
        if layout == "summary":
            closing = [self._compact_text(item, max_len=68) for item in (package.get("closing_points") or [])]
            closing = [item for item in closing if item]
            expanded = self._expand_content_lines([*bullets, *closing, *supplemental], min_count=6, max_count=10)
            return [title, *expanded]
        if layout == "agenda":
            numbered = [f"{i + 1}. {item}" for i, item in enumerate(bullets[:8])]
            return [title, *numbered]

        expanded = self._expand_content_lines([*bullets, *supplemental], min_count=6, max_count=10)
        return [title, *expanded]

    def _fill_template_slide_text(self, slide, text_lines: list[str]):
        title_shape, body_shapes = self._collect_text_shapes(slide)
        if not title_shape and not body_shapes:
            return

        title = (text_lines[0] if text_lines else "")
        body_lines = text_lines[1:] if len(text_lines) > 1 else []

        if title_shape:
            self._set_shape_text(title_shape, [title], role="title")

        if body_shapes:
            body_groups = self._split_lines_by_shape_capacity(body_lines, body_shapes)
            for idx, shape in enumerate(body_shapes, start=0):
                group = body_groups[idx] if idx < len(body_groups) else []
                if group:
                    self._set_shape_text(shape, group, role="body")
                elif getattr(shape, "is_placeholder", False) or self._looks_like_placeholder_text(self._extract_shape_text(shape)):
                    self._set_shape_text(shape, [], role="body")

    def _collect_text_shapes(self, slide):
        min_width = int(Inches(1.8))
        min_height = int(Inches(0.35))
        ignore_narrow_width = int(Inches(1.2))
        ignore_narrow_height = int(Inches(2.4))

        title_candidates = []
        body_candidates = []
        neutral_candidates = []

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            width = getattr(shape, "width", 0)
            height = getattr(shape, "height", 0)
            if width <= 0 or height <= 0:
                continue

            # 过滤装饰性竖排细条文本框，避免出现“文字竖排挤压”问题。
            if width <= ignore_narrow_width and height >= ignore_narrow_height:
                continue

            if width < min_width or height < min_height:
                continue

            current_text = self._extract_shape_text(shape)
            placeholder_kind = self._placeholder_kind(shape)
            if placeholder_kind == "title":
                title_candidates.append(shape)
                continue
            if placeholder_kind == "body":
                body_candidates.append(shape)
                continue
            if self._looks_like_placeholder_text(current_text):
                body_candidates.append(shape)
                continue
            neutral_candidates.append(shape)

        sort_key = lambda shp: (int(getattr(shp, "top", 0)), int(getattr(shp, "left", 0)), -int(getattr(shp, "width", 0) * getattr(shp, "height", 0)))
        title_candidates.sort(key=sort_key)
        body_candidates.sort(key=sort_key)
        neutral_candidates.sort(key=sort_key)

        title_shape = title_candidates[0] if title_candidates else (neutral_candidates[0] if neutral_candidates else None)

        body_shapes = list(body_candidates)
        for shape in neutral_candidates:
            if shape is title_shape:
                continue
            body_shapes.append(shape)

        # 控制正文文本框数量，避免把大量装饰文本都写满。
        body_shapes = body_shapes[:6]
        return title_shape, body_shapes

    def _placeholder_kind(self, shape):
        if not getattr(shape, "is_placeholder", False):
            return "neutral"
        try:
            ph_name = str(getattr(shape.placeholder_format.type, "name", "")).upper()
        except Exception:
            ph_name = ""
        if "TITLE" in ph_name:
            return "title"
        if any(token in ph_name for token in ("BODY", "OBJECT", "CONTENT", "SUBTITLE", "TEXT")):
            return "body"
        return "neutral"

    def _set_shape_text(self, shape, lines: list[str], role: str = "body"):
        normalized = [re.sub(r"\s+", " ", str(line or "")).strip() for line in (lines or [])]
        normalized = [line for line in normalized if line]
        text_frame = shape.text_frame
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        text_frame.vertical_anchor = MSO_ANCHOR.TOP
        text_frame.word_wrap = True
        text_frame.margin_left = int(Inches(0.08))
        text_frame.margin_right = int(Inches(0.08))
        text_frame.margin_top = int(Inches(0.05))
        text_frame.margin_bottom = int(Inches(0.05))
        text_frame.clear()
        if not normalized:
            return

        line_capacity = self._estimate_shape_line_capacity(shape)
        char_capacity = self._estimate_shape_char_capacity(shape)
        safe_lines = []
        current_chars = 0
        for line in normalized:
            line_trim = self._compact_text(line, max_len=char_capacity)
            if len(safe_lines) >= line_capacity:
                break
            if current_chars + len(line_trim) > char_capacity * max(2, line_capacity // 2):
                break
            safe_lines.append(line_trim)
            current_chars += len(line_trim)

        if not safe_lines:
            safe_lines = [self._compact_text(normalized[0], max_len=max(12, char_capacity // 2))]

        is_title = str(role or "body").lower() == "title"
        base_font_size = PptPt(30 if is_title else 20)
        if not is_title:
            shape_height_in = float(getattr(shape, "height", 0)) / 914400.0
            if shape_height_in <= 1.0:
                base_font_size = PptPt(16)
            elif shape_height_in <= 1.5:
                base_font_size = PptPt(18)

        for index, line in enumerate(safe_lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.text = line
            paragraph.level = 0
            paragraph.alignment = PP_ALIGN.LEFT if not is_title else PP_ALIGN.CENTER
            paragraph.line_spacing = 1.18 if not is_title else 1.1
            paragraph.space_before = PptPt(0)
            paragraph.space_after = PptPt(3 if not is_title else 1)
            paragraph.font.bold = bool(is_title)
            paragraph.font.size = base_font_size

    def _split_lines_by_shape_capacity(self, lines: list[str], shapes: list):
        if not shapes:
            return []
        if not lines:
            return [[] for _ in shapes]

        estimated_caps = [self._estimate_shape_line_capacity(shape) for shape in shapes]
        estimated_chars = [self._estimate_shape_char_capacity(shape) for shape in shapes]

        groups = [[] for _ in shapes]
        cursor = 0

        # 第一轮：每个文本框至少分配一条，提升画面填充感。
        for i in range(len(shapes)):
            if cursor >= len(lines):
                break
            groups[i].append(lines[cursor])
            cursor += 1

        # 第二轮：按容量补充分配剩余文本。
        while cursor < len(lines):
            progress = False
            for i, cap in enumerate(estimated_caps):
                if cursor >= len(lines):
                    break
                if len(groups[i]) < cap:
                    groups[i].append(lines[cursor])
                    cursor += 1
                    progress = True
            if not progress:
                break

        # 仍有剩余时，将剩余要点压缩合并到最大正文框末尾，避免直接堆叠导致遮挡或丢失。
        if cursor < len(lines):
            biggest = max(range(len(shapes)), key=lambda idx: int(getattr(shapes[idx], "width", 0)) * int(getattr(shapes[idx], "height", 0)))
            remaining = [self._compact_text(item, max_len=48) for item in lines[cursor:cursor + 5]]
            remaining = [item for item in remaining if item]
            if remaining:
                merged = "；".join(remaining)
                merged = self._compact_text(merged, max_len=max(24, estimated_chars[biggest] * 2))
                groups[biggest].append(merged)

        return groups

    def _estimate_shape_line_capacity(self, shape) -> int:
        height = int(getattr(shape, "height", 0))
        line_height_emu = int(Inches(0.42))
        return max(1, min(9, height // max(line_height_emu, 1)))

    def _estimate_shape_char_capacity(self, shape) -> int:
        width_in = max(1.0, float(getattr(shape, "width", 0)) / 914400.0)
        # 中文投屏场景按保守估计，避免横向溢出。
        return max(16, int(width_in * 9))

    def _extract_shape_text(self, shape) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        parts = []
        for paragraph in shape.text_frame.paragraphs:
            text = re.sub(r"\s+", " ", str(paragraph.text or "")).strip()
            if text:
                parts.append(text)
        return " ".join(parts).strip()

    def _looks_like_placeholder_text(self, text: str) -> bool:
        source = str(text or "").strip().lower()
        if not source:
            return True
        if any(hint in source for hint in self.PLACEHOLDER_HINTS):
            return True
        return False

    def _expand_content_lines(self, lines: list[str], min_count: int = 6, max_count: int = 10) -> list[str]:
        expanded = []
        for line in lines:
            clean = self._compact_text(line, max_len=70)
            if not clean:
                continue
            expanded.append(clean)
            segments = [re.sub(r"\s+", " ", seg).strip() for seg in re.split(r"[；;。]", clean) if re.sub(r"\s+", " ", seg).strip()]
            for seg in segments:
                if len(seg) >= 8 and seg not in expanded:
                    expanded.append(self._compact_text(seg, max_len=62))
            if len(expanded) >= max_count:
                break

        if len(expanded) < min_count and expanded:
            base = list(expanded)
            i = 0
            while len(expanded) < min_count:
                expanded.append(base[i % len(base)])
                i += 1

        dedup = []
        for item in expanded:
            if item not in dedup:
                dedup.append(item)
        return dedup[:max_count]

    def _remove_slide_at(self, prs: Presentation, index: int):
        if index < 0 or index >= len(prs.slides):
            return
        slide_id = prs.slides._sldIdLst[index]
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)

    def _normalize_slides_for_ppt(self, slides: list[dict]):
        max_bullets = {
            "cover": 4,
            "agenda": 6,
            "content": 5,
            "case": 6,
            "interactive": 3,
            "summary": 6,
        }
        normalized = []
        for slide in slides or []:
            layout = str(slide.get("layout") or "content")
            limit = max_bullets.get(layout, 5)
            bullets = []
            for index, raw in enumerate((slide.get("bullets") or [])[:limit]):
                clean = self._compact_text(raw, max_len=44 if index == 0 else 54)
                if clean:
                    bullets.append(clean)
            normalized.append(
                {
                    "title": self._compact_text(slide.get("title") or "未命名页面", max_len=30),
                    "layout": layout,
                    "bullets": bullets,
                }
            )
        return normalized

    def _compact_text(self, value, max_len: int = 54):
        text = re.sub(r"\s+", " ", str(value or "")).strip()
        if len(text) <= max_len:
            return text
        return f"{text[:max_len - 1]}…"

    def build_doc_preview(self, package: dict):
        return {
            "sections": [
                {"title": "教学目标", "items": self._teaching_objectives(package)},
                {"title": "教学方法", "items": self._teaching_methods(package)},
                {"title": "课堂活动设计", "items": self._classroom_activities(package)},
                {"title": "课后作业", "items": self._homework_design(package)},
            ]
        }

    def build_ppt_slide_thumbnails(self, ppt_filename: str, max_slides: int = 20):
        ppt_path = self.output_dir / str(ppt_filename or "").strip()
        if not ppt_path.exists() or ppt_path.suffix.lower() != ".pptx":
            return []

        stat = ppt_path.stat()
        digest = hashlib.md5(f"{ppt_path.resolve()}::{stat.st_mtime_ns}::{stat.st_size}".encode("utf-8")).hexdigest()[:16]
        max_count = max(1, min(int(max_slides or 20), 40))

        urls = []
        for index in range(1, max_count + 1):
            target_name = f"ppt_{digest}_{index:03d}.png"
            target_path = self.preview_dir / target_name
            if target_path.exists():
                urls.append(f"/static/generated_ppt_previews/{quote(target_name)}")
            else:
                break
        if urls:
            return urls

        export_source_path = ppt_path
        temp_source_path = None
        if any(ord(ch) > 127 for ch in str(ppt_path)):
            temp_source_path = self.preview_dir / f"thumb_src_{digest}.pptx"
            try:
                shutil.copy2(ppt_path, temp_source_path)
                export_source_path = temp_source_path
            except Exception:
                export_source_path = ppt_path

        src = str(export_source_path.resolve())
        out_dir = str(self.preview_dir.resolve())
        if self._export_thumbnails_via_win32com(Path(src), digest=digest, max_count=max_count):
            urls = []
            for index in range(1, max_count + 1):
                target_name = f"ppt_{digest}_{index:03d}.png"
                target_path = self.preview_dir / target_name
                if not target_path.exists():
                    break
                urls.append(f"/static/generated_ppt_previews/{quote(target_name)}")
            if urls:
                if temp_source_path and temp_source_path.exists():
                    try:
                        temp_source_path.unlink()
                    except OSError:
                        pass
                return urls

        src_ps = src.replace("'", "''")
        out_dir_ps = out_dir.replace("'", "''")
        script = (
            "$ErrorActionPreference='Stop';"
            "$ppt=$null; $pres=$null;"
            "try {"
            "$ppt=New-Object -ComObject PowerPoint.Application;"
            f"$pres=$ppt.Presentations.Open('{src_ps}', $false, $true, $false);"
            f"$max=[Math]::Min($pres.Slides.Count,{max_count});"
            "for($i=1; $i -le $max; $i++){"
            f"$name=('ppt_{digest}_' + $i.ToString('000') + '.png');"
            f"$target=Join-Path '{out_dir_ps}' $name;"
            "$pres.Slides.Item($i).Export($target,'PNG',960,540);"
            "}"
            "Write-Output $max;"
            "} finally {"
            "if($pres -ne $null){try{$pres.Close()}catch{}}"
            "if($ppt -ne $null){try{$ppt.Quit()}catch{}}"
            "}"
        )
        try:
            result = subprocess.run(
                ["powershell", "-NoProfile", "-STA", "-Command", script],
                capture_output=True,
                text=True,
                timeout=60,
            )
        except Exception:
            return []
        finally:
            if temp_source_path and temp_source_path.exists():
                try:
                    temp_source_path.unlink()
                except OSError:
                    pass

        urls = []
        for index in range(1, max_count + 1):
            target_name = f"ppt_{digest}_{index:03d}.png"
            target_path = self.preview_dir / target_name
            if not target_path.exists():
                break
            urls.append(f"/static/generated_ppt_previews/{quote(target_name)}")
        return urls

    def _export_thumbnails_via_win32com(self, source_path: Path, digest: str, max_count: int) -> bool:
        try:
            import pythoncom  # type: ignore
            import win32com.client  # type: ignore
        except Exception:
            return False

        app = None
        presentation = None
        initialized = False
        try:
            pythoncom.CoInitialize()
            initialized = True
            app = win32com.client.DispatchEx("PowerPoint.Application")
            presentation = app.Presentations.Open(str(source_path.resolve()), False, True, False)
            total = min(int(presentation.Slides.Count), int(max_count))
            if total <= 0:
                return False
            for index in range(1, total + 1):
                target_name = f"ppt_{digest}_{index:03d}.png"
                target_path = self.preview_dir / target_name
                presentation.Slides.Item(index).Export(str(target_path.resolve()), "PNG", 960, 540)
            return True
        except Exception:
            return False
        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            if initialized:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

    def _resolve_blank_layout(self, prs: Presentation):
        for layout in prs.slide_layouts:
            if "blank" in str(layout.name).lower() or "空白" in str(layout.name):
                return layout
        return prs.slide_layouts[len(prs.slide_layouts) - 1]

    def _add_blank_slide(self, prs: Presentation):
        return prs.slides.add_slide(self._resolve_blank_layout(prs))

    def _build_doc(self, file_path: Path, package: dict):
        self._apply_theme(package.get("theme_style") or package.get("summary", {}).get("style") or "教育蓝")
        document = Document()
        document.styles["Normal"].font.name = "Microsoft YaHei"
        document.styles["Normal"].font.size = Pt(11)

        heading = document.add_heading(package["title"], level=0)
        heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        heading.runs[0].font.color.rgb = DocRGBColor(*self.theme["accent"])

        document.add_heading("一、教学概览", level=1)
        self._write_kv(document, "课程主题", package["summary"].get("course_theme", ""))
        self._write_kv(document, "知识点", package["summary"].get("knowledge_points", ""))
        self._write_kv(document, "重难点", package["summary"].get("key_difficulties", ""))
        self._write_kv(document, "课时安排", package["summary"].get("lesson_periods", ""))
        self._write_kv(document, "课件模板", package.get("theme_style", ""))

        document.add_heading("二、教学目标", level=1)
        for item in self._teaching_objectives(package):
            document.add_paragraph(item, style="List Bullet")

        document.add_heading("三、教学方法", level=1)
        for item in self._teaching_methods(package):
            document.add_paragraph(item, style="List Bullet")

        document.add_heading("四、课堂活动设计", level=1)
        for item in self._classroom_activities(package):
            document.add_paragraph(item, style="List Bullet")

        document.add_heading("五、教学流程", level=1)
        for index, slide in enumerate(package["slides"], start=1):
            document.add_heading(f"{index}. {slide['title']}", level=2)
            for bullet in slide.get("bullets", []):
                document.add_paragraph(bullet, style="List Bullet")

        document.add_heading("六、课后作业", level=1)
        for item in self._homework_design(package):
            document.add_paragraph(item, style="List Bullet")

        document.add_heading("七、课堂收束", level=1)
        for point in package.get("closing_points", []):
            document.add_paragraph(point, style="List Bullet")

        document.add_heading("八、资料融入", level=1)
        for item in package.get("references", []):
            document.add_paragraph(item, style="List Bullet")

        document.save(str(file_path))

    def _apply_theme(self, style_name: str):
        style_name = str(style_name or "")
        for key, value in self.THEMES.items():
            if key in style_name:
                self.theme = value
                return
        external_style_map = {
            "math": "学术风",
            "function": "学术风",
            "ratio": "学术风",
            "angle": "几何风",
            "geometric": "几何风",
            "blackboard": "学术风",
            "cute": "可爱风",
            "playful": "卡通风",
            "notebook": "手账风",
            "matisse": "插画风",
            "history": "中国风",
            "medical": "商务风",
            "psychology": "莫兰迪风",
            "stem": "科技风",
        }
        lower_style = style_name.lower()
        for keyword, mapped_theme in external_style_map.items():
            if keyword in lower_style and mapped_theme in self.THEMES:
                self.theme = self.THEMES[mapped_theme]
                return
        self.theme = self.THEMES["教育蓝"]

    def _set_background(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*color)

    def _text_box(self, slide, left, top, width, height, text, size, color, bold=False):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.font.size = PptPt(size)
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.font.color.rgb = RGBColor(*color)
        return box

    def _add_top_band(self, slide):
        band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.2))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(*self.theme["accent"])
        band.line.fill.background()

    def _add_cover_slide(self, prs: Presentation, package: dict):
        slide = self._add_blank_slide(prs)
        self._set_background(slide, self.theme["bg"])
        self._add_top_band(slide)
        family = self.theme["family"]

        if family in {"neon", "dark"}:
            panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.0), Inches(7.6), Inches(5.4))
            panel.fill.solid()
            panel.fill.fore_color.rgb = RGBColor(*self.theme["soft"])
            panel.line.fill.background()
            accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.85), Inches(0.95), Inches(3.4), Inches(5.6))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(*self.theme["accent"])
            accent.fill.transparency = 0.18
            accent.line.fill.background()
            title_color = self.theme["text"]
        elif family in {"playful", "airy"}:
            accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.1), Inches(11.4), Inches(5.2))
            accent.fill.solid()
            accent.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            accent.line.color.rgb = RGBColor(*self.theme["accent"])
            for shape_left in (9.8, 10.9, 11.8):
                blob = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(shape_left), Inches(0.75 + (shape_left - 9.8) * 0.2), Inches(0.65), Inches(0.65))
                blob.fill.solid()
                blob.fill.fore_color.rgb = RGBColor(*self.theme["accent"])
                blob.fill.transparency = 0.2
                blob.line.fill.background()
            title_color = self.theme["dark_text"]
        else:
            left_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.9), Inches(0.8), Inches(3.1), Inches(5.9))
            left_panel.fill.solid()
            left_panel.fill.fore_color.rgb = RGBColor(*self.theme["accent"])
            left_panel.fill.transparency = 0.08
            left_panel.line.fill.background()
            body_panel = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(3.35), Inches(1.0), Inches(8.8), Inches(5.35))
            body_panel.fill.solid()
            body_panel.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            body_panel.line.fill.background()
            title_color = self.theme["dark_text"]

        self._text_box(slide, 1.05, 1.25, 10.0, 1.45, package["title"], 28, title_color, bold=True)
        items = [
            f"课程主题：{package['summary'].get('course_theme', '')}",
            f"知识点：{package['summary'].get('knowledge_points', '')}",
            f"课时安排：{package['summary'].get('lesson_periods', '')}",
            f"课件模板：{package.get('theme_style', '')}",
        ]
        info = self._text_box(slide, 1.08, 3.0, 7.2, 2.6, "", 18, title_color)
        for idx, item in enumerate(items):
            p = info.text_frame.paragraphs[0] if idx == 0 else info.text_frame.add_paragraph()
            p.text = item
            p.font.size = PptPt(17)
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(*title_color)

    def _add_agenda_slide(self, prs: Presentation, slide_data: dict):
        slide = self._add_blank_slide(prs)
        self._set_background(slide, self.theme["soft"])
        self._add_top_band(slide)
        self._text_box(slide, 0.8, 0.55, 6.5, 0.8, slide_data["title"], 24, self.theme["dark_text"], bold=True)

        family = self.theme["family"]
        bullets = slide_data.get("bullets", [])[:6]
        for index, bullet in enumerate(bullets):
            if family in {"editorial", "split"}:
                left = 0.9 + (index % 3) * 4.0
                top = 1.6 + (index // 3) * 2.1
                card_w = 3.5
                card_h = 1.75
            else:
                left = 0.95 + (index % 2) * 6.05
                top = 1.6 + (index // 2) * 1.45
                card_w = 5.35
                card_h = 1.0
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(card_h))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            card.line.color.rgb = RGBColor(*self.theme["accent"])
            card.text_frame.word_wrap = True
            p = card.text_frame.paragraphs[0]
            p.text = f"{index + 1:02d}  {bullet}"
            p.font.size = PptPt(18 if card_h < 1.2 else 16)
            p.font.bold = True
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(*self.theme["dark_text"])

    def _add_content_slide(self, prs: Presentation, slide_data: dict):
        slide = self._add_blank_slide(prs)
        self._set_background(slide, self.theme["soft"])
        self._add_top_band(slide)
        self._text_box(slide, 0.85, 0.55, 7.0, 0.8, slide_data["title"], 24, self.theme["dark_text"], bold=True)
        family = self.theme["family"]

        if family in {"geometry", "flat", "glass"}:
            left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.45), Inches(4.0), Inches(4.9))
            left.fill.solid()
            left.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            left.line.color.rgb = RGBColor(*self.theme["accent"])
            right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(1.45), Inches(7.15), Inches(4.9))
            right.fill.solid()
            right.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            right.line.fill.background()
            targets = [(slide_data.get("bullets", [])[:2], left), (slide_data.get("bullets", [])[2:5], right)]
            for bullets, box in targets:
                frame = box.text_frame
                frame.word_wrap = True
                for idx, bullet in enumerate(bullets):
                    p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
                    p.text = bullet
                    p.font.size = PptPt(20 if idx == 0 else 17)
                    p.font.bold = idx == 0
                    p.font.name = "Microsoft YaHei"
                    p.font.color.rgb = RGBColor(*self.theme["dark_text"])
        elif family in {"playful", "airy"}:
            banner = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.5), Inches(11.55), Inches(4.8))
            banner.fill.solid()
            banner.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            banner.line.color.rgb = RGBColor(*self.theme["accent"])
            for idx, bullet in enumerate(slide_data.get("bullets", [])[:5]):
                p = banner.text_frame.paragraphs[0] if idx == 0 else banner.text_frame.add_paragraph()
                p.text = bullet
                p.font.size = PptPt(21 if idx == 0 else 18)
                p.font.bold = idx == 0
                p.font.name = "Microsoft YaHei"
                p.font.color.rgb = RGBColor(*self.theme["dark_text"])
        else:
            side = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.7), Inches(1.45), Inches(0.22), Inches(5.2))
            side.fill.solid()
            side.fill.fore_color.rgb = RGBColor(*self.theme["accent"])
            side.line.fill.background()
            content = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.55), Inches(11.4), Inches(4.95))
            content.fill.solid()
            content.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            content.line.fill.background()
            content.text_frame.word_wrap = True
            for idx, bullet in enumerate(slide_data.get("bullets", [])[:5]):
                p = content.text_frame.paragraphs[0] if idx == 0 else content.text_frame.add_paragraph()
                p.text = bullet
                p.font.size = PptPt(22 if idx == 0 else 18)
                p.font.bold = idx == 0
                p.font.name = "Microsoft YaHei"
                p.font.color.rgb = RGBColor(*self.theme["dark_text"])

    def _add_case_slide(self, prs: Presentation, slide_data: dict):
        slide = self._add_blank_slide(prs)
        self._set_background(slide, self.theme["soft"])
        self._add_top_band(slide)
        self._text_box(slide, 0.85, 0.55, 7.0, 0.8, slide_data["title"], 24, self.theme["dark_text"], bold=True)

        left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95), Inches(1.65), Inches(3.8), Inches(4.95))
        left.fill.solid()
        left.fill.fore_color.rgb = RGBColor(*self.theme["accent"])
        left.fill.transparency = 0.12 if self.theme["family"] not in {"dark", "neon"} else 0.04
        left.line.fill.background()
        head = left.text_frame.paragraphs[0]
        head.text = "案例焦点"
        head.font.size = PptPt(23)
        head.font.bold = True
        head.font.name = "Microsoft YaHei"
        head.font.color.rgb = RGBColor(*self.theme["dark_text"])
        for bullet in slide_data.get("bullets", [])[:2]:
            p = left.text_frame.add_paragraph()
            p.text = bullet
            p.font.size = PptPt(17)
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(*self.theme["dark_text"])

        right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.0), Inches(1.65), Inches(7.35), Inches(4.95))
        right.fill.solid()
        right.fill.fore_color.rgb = RGBColor(*self.theme["card"])
        right.line.color.rgb = RGBColor(*self.theme["accent"])
        right.text_frame.word_wrap = True
        for idx, bullet in enumerate(slide_data.get("bullets", [])[2:6]):
            p = right.text_frame.paragraphs[0] if idx == 0 else right.text_frame.add_paragraph()
            p.text = bullet
            p.font.size = PptPt(18)
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(*self.theme["dark_text"])

    def _add_interactive_slide(self, prs: Presentation, slide_data: dict):
        slide = self._add_blank_slide(prs)
        self._set_background(slide, self.theme["soft"])
        self._add_top_band(slide)
        self._text_box(slide, 0.85, 0.55, 7.0, 0.8, slide_data["title"], 24, self.theme["dark_text"], bold=True)

        for idx, bullet in enumerate(slide_data.get("bullets", [])[:3]):
            left = 0.95 + idx * 4.05
            card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(3.45), Inches(3.95))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(*self.theme["card"])
            card.line.color.rgb = RGBColor(*self.theme["accent"])
            badge = card.text_frame.paragraphs[0]
            badge.text = f"互动 {idx + 1}"
            badge.font.size = PptPt(18)
            badge.font.bold = True
            badge.font.name = "Microsoft YaHei"
            badge.font.color.rgb = RGBColor(*self.theme["accent"])
            body = card.text_frame.add_paragraph()
            body.text = bullet
            body.font.size = PptPt(16)
            body.font.name = "Microsoft YaHei"
            body.font.color.rgb = RGBColor(*self.theme["dark_text"])

    def _add_summary_slide(self, prs: Presentation, slide_data: dict, package: dict):
        slide = self._add_blank_slide(prs)
        self._set_background(slide, self.theme["bg"])
        self._add_top_band(slide)
        self._text_box(slide, 0.85, 0.55, 7.0, 0.8, slide_data["title"], 24, self.theme["text"], bold=True)

        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(1.65), Inches(11.95), Inches(4.95))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*self.theme["card"])
        box.line.fill.background()

        items = []
        for item in slide_data.get("bullets", []) + package.get("closing_points", []):
            if item not in items:
                items.append(item)
        for idx, bullet in enumerate(items[:6]):
            p = box.text_frame.paragraphs[0] if idx == 0 else box.text_frame.add_paragraph()
            p.text = bullet
            p.font.size = PptPt(19)
            p.font.name = "Microsoft YaHei"
            p.font.color.rgb = RGBColor(*self.theme["dark_text"])

    def _write_kv(self, document: Document, label: str, value: str):
        paragraph = document.add_paragraph()
        run = paragraph.add_run(f"{label}：")
        run.bold = True
        run.font.color.rgb = DocRGBColor(*self.theme["accent"])
        paragraph.add_run(value)

    def _teaching_objectives(self, package: dict):
        summary = package.get("summary") or {}
        knowledge = str(summary.get("knowledge_points") or "核心知识点").strip()
        difficulties = str(summary.get("key_difficulties") or "重难点").strip()
        return [
            f"知识与技能：理解并掌握“{knowledge}”的核心概念、关键方法与典型应用。",
            f"过程与方法：通过问题驱动、案例分析与分层练习，突破“{difficulties}”。",
            "情感与价值：培养严谨推理意识与自主探究习惯，提升课堂参与度与表达能力。",
        ]

    def _teaching_methods(self, package: dict):
        slides = package.get("slides", [])
        layouts = {str(slide.get("layout") or "") for slide in slides}
        methods = ["讲授法：用于知识框架搭建与关键概念讲解。", "启发式问答：通过追问引导学生主动表达思路。"]
        if "case" in layouts:
            methods.append("案例教学法：借助典型例题/情境案例落实知识迁移。")
        if "interactive" in layouts:
            methods.append("互动活动法：通过课堂互动或小游戏强化即时反馈。")
        methods.append("分层练习法：按照基础—提升两层任务组织练习与评价。")
        return methods

    def _classroom_activities(self, package: dict):
        summary = package.get("summary") or {}
        theme = str(summary.get("course_theme") or "本课主题").strip()
        activities = [
            f"导入活动：用与“{theme}”相关的问题情境唤醒前置知识（约5分钟）。",
            "建构活动：围绕核心概念进行板书/结构图梳理，明确方法步骤（约15分钟）。",
            "练习活动：组织1-2个分层任务，强调方法比较与易错点纠偏（约15分钟）。",
        ]
        for slide in package.get("slides", []):
            if slide.get("layout") == "interactive":
                first = next((str(item).strip() for item in slide.get("bullets", []) if str(item).strip()), "课堂互动任务")
                activities.append(f"互动活动：{first}（约5分钟）。")
                break
        activities.append("小结活动：通过出口任务或口头复述检验学习达成（约5分钟）。")
        return activities

    def _homework_design(self, package: dict):
        summary = package.get("summary") or {}
        knowledge = str(summary.get("knowledge_points") or "本课知识点").strip()
        difficulties = str(summary.get("key_difficulties") or "重难点").strip()
        return [
            f"基础巩固：围绕“{knowledge}”完成基础题 3-5 题，重点核对步骤规范。",
            f"能力提升：针对“{difficulties}”设计1-2题综合应用题，要求写出解题思路。",
            "拓展任务：结合真实情境自拟一个相关问题，并给出分析或求解过程。",
        ]
