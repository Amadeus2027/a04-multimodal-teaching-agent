from __future__ import annotations

import json
import re
import ast
import copy
import types
import base64
from pathlib import Path
import hashlib
import zipfile
import subprocess
import os
from html import unescape
from urllib.parse import quote, quote_plus, urljoin, urlparse, parse_qs, unquote
from urllib.request import Request, urlopen
from concurrent.futures import ThreadPoolExecutor, as_completed

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches, Pt

from services.document_service import DocumentExportService
from services.interactive_service import InteractiveService
from services.session_manager import REQUIRED_FIELDS
from services.video_renderer import VideoRenderer


class ContentService:
    """Coordinates clarification, retrieval, instruction fusion, and package generation."""

    CREATIVE_KEYWORDS = ["动画", "动效", "小游戏", "互动游戏", "闯关", "拖拽", "配对", "课堂互动", "创意", "演示视频", "知识演示", "动态演示", "可视化", "geogebra", "desmos", "math3d", "wolfram", "链接页"]
    OPEN_DEMO_LINKS = [
        ("GeoGebra", "https://www.geogebra.org/"),
        ("Desmos Calculator", "https://www.desmos.com/calculator"),
        ("Math3D", "https://www.math3d.org/"),
        ("Wolfram Demonstrations", "https://demonstrations.wolfram.com/"),
    ]
    OPEN_DEMO_SITES = [
        {
            "name": "GeoGebra",
            "url": "https://www.geogebra.org/",
            "domain": "geogebra.org",
            "search": "https://www.geogebra.org/search/{query}",
            "search_templates": [
                "https://www.geogebra.org/search/{query}",
                "https://www.geogebra.org/materials?query={query}",
                "https://duckduckgo.com/html/?q=site%3Ageogebra.org+{query_plus}",
                "https://www.bing.com/search?q=site%3Ageogebra.org+{query_plus}",
            ],
            "detail_hints": ["/m/", "/material/", "/materials/", "activity"],
            "avoid_hints": ["/search", "login", "signup", "privacy", "terms"],
            "default_summary": "覆盖函数图像、几何构造与课堂活动，适合做概念可视化与探究任务。",
            "suggestions": ["函数图像动态演示", "参数变化观察", "课堂互动活动"],
            "brief_intro": "开源动态数学软件，支持函数绘图、几何作图与代数运算，可创建交互式课堂活动，适合概念可视化与探究式教学。",
        },
        {
            "name": "Desmos Calculator",
            "url": "https://www.desmos.com/calculator",
            "domain": "desmos.com",
            "search": "https://teacher.desmos.com/search?q={query}",
            "search_templates": [
                "https://teacher.desmos.com/search?q={query}",
                "https://duckduckgo.com/html/?q=site%3Adesmos.com+{query_plus}",
                "https://www.bing.com/search?q=site%3Adesmos.com+{query_plus}",
            ],
            "detail_hints": ["/calculator/", "/3d/", "/activitybuilder/", "/activities/"],
            "avoid_hints": ["/search", "login", "signup", "privacy", "terms"],
            "default_summary": "图形计算器响应快，适合展示函数、极限邻域变化与参数联动。",
            "suggestions": ["函数与导数图像", "极限邻域放大", "参数滑块演示"],
            "brief_intro": "在线图形计算器，支持函数、参数方程与数据可视化，响应迅速，适合课堂实时演示函数变化与极限邻域。",
        },
        {
            "name": "Math3D",
            "url": "https://www.math3d.org/",
            "domain": "math3d.org",
            "search": "https://www.math3d.org/",
            "search_templates": [
                "https://duckduckgo.com/html/?q=site%3Amath3d.org+{query_plus}",
                "https://www.bing.com/search?q=site%3Amath3d.org+{query_plus}",
            ],
            "detail_hints": ["/examples", "/example", "/gallery", "/lesson", "/surface"],
            "avoid_hints": ["/search", "login", "signup", "privacy", "terms"],
            "default_summary": "面向三维函数与曲面演示，可用于空间几何与多元函数直观讲解。",
            "suggestions": ["三维曲面观察", "截面变化", "空间参数影响"],
            "brief_intro": "在线三维数学可视化工具，支持曲面、空间曲线与向量场演示，适合多元函数与空间几何直观讲解。",
        },
        {
            "name": "Wolfram Demonstrations",
            "url": "https://demonstrations.wolfram.com/",
            "domain": "demonstrations.wolfram.com",
            "search": "https://demonstrations.wolfram.com/search.html?query={query}",
            "search_templates": [
                "https://demonstrations.wolfram.com/search.html?query={query}",
                "https://duckduckgo.com/html/?q=site%3Ademonstrations.wolfram.com+{query_plus}",
                "https://www.bing.com/search?q=site%3Ademonstrations.wolfram.com+{query_plus}",
            ],
            "detail_hints": ["/demonstrations/", "/topics/", "/detail"],
            "avoid_hints": ["/search", "login", "signup", "privacy", "terms"],
            "default_summary": "汇集大量数学与科学交互演示，适合做拓展案例与课后探索。",
            "suggestions": ["可视化案例拓展", "跨学科演示", "探究式作业素材"],
            "brief_intro": "Wolfram 演示项目收录大量交互式数学与科学演示，涵盖微积分、线性代数等领域，适合拓展案例与课后探究。",
        },
    ]

    TEMPLATE_PROFILES = [
        {"name": "学术风", "scene": "概念讲解、推理证明、考试复习", "tags": ["学术", "严谨", "证明", "推理", "理论", "复习"]},
        {"name": "教育蓝", "scene": "通用课堂、教研公开课、教师培训", "tags": ["通用", "教育", "课堂", "公开课", "培训"]},
        {"name": "极简风", "scene": "信息密度高、结构化讲解", "tags": ["极简", "简洁", "条理", "结构"]},
        {"name": "科技风", "scene": "信息技术、科学、编程、AI", "tags": ["科技", "数字", "信息", "编程", "人工智能", "创新"]},
        {"name": "卡通风", "scene": "低龄课堂、趣味导入、闯关互动", "tags": ["卡通", "趣味", "游戏", "闯关", "低龄"]},
        {"name": "可爱风", "scene": "儿童化表达、启蒙课堂", "tags": ["可爱", "儿童", "启蒙", "亲和", "活泼"]},
        {"name": "插画风", "scene": "故事化课程、绘本式表达", "tags": ["插画", "故事", "绘本", "情境"]},
        {"name": "清新风", "scene": "语文、英语、综合素养课", "tags": ["清新", "轻盈", "自然", "柔和"]},
        {"name": "中国风", "scene": "传统文化、历史、古诗文", "tags": ["中国风", "传统文化", "古典", "历史", "诗词"]},
        {"name": "新中式", "scene": "传统文化、校本课程、高级感展示", "tags": ["新中式", "国风", "传统", "文化", "典雅"]},
        {"name": "商务风", "scene": "汇报课、成果展示、说课", "tags": ["商务", "汇报", "成果", "展示", "说课"]},
        {"name": "杂志风", "scene": "案例展示、作品讲评、项目化学习", "tags": ["杂志", "案例", "作品", "项目", "展示"]},
        {"name": "几何风", "scene": "数学、图形、逻辑结构", "tags": ["几何", "数学", "图形", "结构", "逻辑"]},
        {"name": "莫兰迪风", "scene": "温和叙述、心理健康、审美课程", "tags": ["莫兰迪", "柔和", "审美", "心理", "温和"]},
        {"name": "粉彩风", "scene": "低刺激、温暖陪伴型课堂", "tags": ["粉彩", "温暖", "柔和", "陪伴"]},
        {"name": "暖阳风", "scene": "德育、心理、班会、启发式课程", "tags": ["暖阳", "启发", "心理", "德育", "班会"]},
        {"name": "自然风", "scene": "科学观察、生物、劳动教育", "tags": ["自然", "观察", "科学", "生物", "劳动"]},
        {"name": "森林风", "scene": "生态主题、生命教育", "tags": ["森林", "生态", "植物", "生命", "自然"]},
        {"name": "海洋风", "scene": "地理、海洋、环保主题", "tags": ["海洋", "蓝色", "环保", "地理", "水"]},
        {"name": "复古风", "scene": "历史回顾、经典案例", "tags": ["复古", "历史", "经典", "年代"]},
        {"name": "手账风", "scene": "过程记录、学习单式课堂", "tags": ["手账", "记录", "学习单", "任务"]},
        {"name": "未来感", "scene": "前沿主题、科创展示", "tags": ["未来", "前沿", "创新", "科创"]},
        {"name": "黑金风", "scene": "高完成度展示、成果汇报", "tags": ["高端", "汇报", "展示", "黑金"]},
        {"name": "暗夜风", "scene": "深色屏显、沉浸式展示", "tags": ["暗夜", "深色", "沉浸", "夜间"]},
        {"name": "霓虹风", "scene": "未来科技、活动舞台感", "tags": ["霓虹", "未来", "舞台", "冲击"]},
        {"name": "玻璃拟态", "scene": "现代界面感、科技展示", "tags": ["玻璃", "拟态", "现代", "UI"]},
        {"name": "扁平风", "scene": "清晰易读、适合大部分课堂投屏", "tags": ["扁平", "清晰", "通用", "直观"]},
        {"name": "渐变风", "scene": "品牌感、现代感表达", "tags": ["渐变", "现代", "品牌", "流动"]},
    ]

    def __init__(self, llm_client, rag_service, output_dir: Path):
        self.llm_client = llm_client
        self.rag_service = rag_service
        self.export_service = DocumentExportService(output_dir=output_dir)
        self._ensure_interactive_movie_embedding_hook()
        self.interactive_service = InteractiveService(llm_client=llm_client)
        self.video_renderer = VideoRenderer(timeout_seconds=int(os.getenv("INTERACTIVE_RENDER_TIMEOUT_SECONDS", "45")))
        self.interactive_assets_root = output_dir / "interactive_assets"
        self.interactive_assets_root.mkdir(parents=True, exist_ok=True)
        self.interactive_clip_duration = float(os.getenv("INTERACTIVE_CLIP_DURATION_SECONDS", "18"))
        self.interactive_max_videos = max(0, int(os.getenv("INTERACTIVE_MAX_VIDEOS", "1")))
        self._site_reference_cache = {}
        self._topic_term_cache = {}
        self.workspace_root = output_dir.parent
        self.template_thumb_dir = self.workspace_root / "static" / "template_thumbs"
        self.template_thumb_dir.mkdir(parents=True, exist_ok=True)
        self.enable_com_thumbnail_fallback = str(os.getenv("ENABLE_TEMPLATE_COM_THUMBNAIL", "0")).strip().lower() in {"1", "true", "yes", "on"}
        self.only_ppt_format_templates = str(os.getenv("ONLY_PPT_FORMAT_TEMPLATES", "1")).strip().lower() in {"1", "true", "yes", "on"}
        external_profiles = self._load_external_template_profiles(output_dir.parent / "ppt_format")
        if self.only_ppt_format_templates and external_profiles:
            self.template_profiles = list(external_profiles)
        else:
            self.template_profiles = list(self.TEMPLATE_PROFILES)
            self.template_profiles.extend(external_profiles)
        dedup = {}
        for item in self.template_profiles:
            dedup[item["name"]] = item
        self.template_profiles = list(dedup.values())
        self.template_names = [item["name"] for item in self.template_profiles]
        self.template_scenes = {item["name"]: item["scene"] for item in self.template_profiles}

    def _update_progress(self, session: dict, percent: int, label: str):
        state = session.get("state") or {}
        state["progress_percent"] = max(0, min(100, int(percent)))
        state["progress_label"] = str(label or "").strip()

    def _load_external_template_profiles(self, template_dir: Path):
        if not template_dir.exists():
            return []
        profiles = []
        for path in sorted(template_dir.glob("*.pptx")):
            name = path.stem.strip()
            if not name:
                continue
            thumbnail_url = self._extract_template_thumbnail(path)
            if not thumbnail_url:
                thumbnail_url = self._build_svg_thumbnail(name=name, scene=self._infer_template_scene(name), key_seed=str(path))
            profiles.append(
                {
                    "name": name,
                    "scene": self._infer_template_scene(name),
                    "tags": self._infer_template_tags(name),
                    "thumbnail_url": thumbnail_url,
                    "source": "external",
                }
            )
        return profiles

    def _build_svg_thumbnail(self, name: str, scene: str, key_seed: str):
        seed = hashlib.md5(str(key_seed).encode("utf-8")).hexdigest()
        bg1 = f"#{seed[:6]}"
        bg2 = f"#{seed[6:12]}"
        txt = (name or "模板")[:28]
        sub = (scene or "课堂模板")[:32]
        safe_txt = txt.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        safe_sub = sub.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
        content = (
            "<svg xmlns='http://www.w3.org/2000/svg' width='960' height='540' viewBox='0 0 960 540'>"
            "<defs>"
            f"<linearGradient id='g' x1='0' y1='0' x2='1' y2='1'><stop offset='0%' stop-color='{bg1}'/><stop offset='100%' stop-color='{bg2}'/></linearGradient>"
            "</defs>"
            "<rect width='960' height='540' fill='url(#g)'/>"
            "<rect x='44' y='44' width='872' height='452' rx='24' fill='rgba(255,255,255,0.84)'/>"
            f"<text x='74' y='170' font-size='42' font-family='Microsoft YaHei, Segoe UI, sans-serif' fill='#1f2937'>{safe_txt}</text>"
            f"<text x='74' y='228' font-size='24' font-family='Microsoft YaHei, Segoe UI, sans-serif' fill='#475569'>{safe_sub}</text>"
            "<text x='74' y='448' font-size='20' font-family='Microsoft YaHei, Segoe UI, sans-serif' fill='#64748b'>A04 Demo Template Preview</text>"
            "</svg>"
        )
        file_name = f"svg_{seed}.svg"
        out_path = self.template_thumb_dir / file_name
        if not out_path.exists():
            out_path.write_text(content, encoding="utf-8")
        return f"/static/template_thumbs/{quote(file_name)}"

    def _extract_template_thumbnail(self, template_path: Path):
        try:
            stat = template_path.stat()
        except OSError:
            return ""

        key = f"{template_path.resolve()}::{stat.st_mtime_ns}::{stat.st_size}"
        digest = hashlib.md5(key.encode("utf-8")).hexdigest()

        # 优先使用模板首页真实导出图，确保模板选择区看到的是“真实模板效果”。
        if self.enable_com_thumbnail_fallback:
            exported_cover = self._export_thumbnail_via_powerpoint(template_path, f"{digest}_cover")
            if exported_cover:
                return exported_cover

        candidates = (
            ("docProps/thumbnail.jpeg", ".jpg"),
            ("docProps/thumbnail.jpg", ".jpg"),
            ("docProps/thumbnail.png", ".png"),
        )
        try:
            with zipfile.ZipFile(template_path, "r") as zf:
                for inner_path, ext in candidates:
                    if inner_path not in zf.namelist():
                        continue
                    target_name = f"{digest}{ext}"
                    target_path = self.template_thumb_dir / target_name
                    if not target_path.exists():
                        target_path.write_bytes(zf.read(inner_path))
                    return f"/static/template_thumbs/{target_name}"
        except Exception:
            return ""
        if not self.enable_com_thumbnail_fallback:
            return ""
        fallback_url = self._export_thumbnail_via_powerpoint(template_path, digest)
        return fallback_url or ""

    def _export_thumbnail_via_powerpoint(self, template_path: Path, digest: str):
        target_name = f"{digest}.png"
        target_path = self.template_thumb_dir / target_name
        if target_path.exists():
            return f"/static/template_thumbs/{target_name}"

        src = str(template_path.resolve())
        dst = str(target_path.resolve())

        if self._export_template_thumbnail_via_win32com(template_path=Path(src), target_path=Path(dst)):
            if target_path.exists():
                return f"/static/template_thumbs/{target_name}"

        src_ps = src.replace("'", "''")
        dst_ps = dst.replace("'", "''")
        script = (
            "$ErrorActionPreference='Stop';"
            "$ppt=$null; $pres=$null;"
            "try {"
            "$ppt=New-Object -ComObject PowerPoint.Application;"
            f"$pres=$ppt.Presentations.Open('{src_ps}', $false, $true, $false);"
            "if($pres.Slides.Count -gt 0){"
            f"$pres.Slides.Item(1).Export('{dst_ps}','PNG',960,540);"
            "}"
            "} finally {"
            "if($pres -ne $null){try{$pres.Close()}catch{}}"
            "if($ppt -ne $null){try{$ppt.Quit()}catch{}}"
            "}"
        )
        try:
            result = subprocess.run(
                ["powershell", "-NoProfile", "-STA", "-Command", script],
                capture_output=True,
                text=True,
                timeout=20,
            )
            if result.returncode == 0 and target_path.exists():
                return f"/static/template_thumbs/{target_name}"
        except Exception:
            return ""
        return ""

    def _export_template_thumbnail_via_win32com(self, template_path: Path, target_path: Path) -> bool:
        try:
            import pythoncom  # type: ignore
            import win32com.client  # type: ignore
        except Exception:
            return False

        app = None
        presentation = None
        initialized = False
        try:
            pythoncom.CoInitialize()
            initialized = True
            app = win32com.client.DispatchEx("PowerPoint.Application")
            presentation = app.Presentations.Open(str(template_path.resolve()), False, True, False)
            if int(presentation.Slides.Count) <= 0:
                return False
            presentation.Slides.Item(1).Export(str(target_path.resolve()), "PNG", 960, 540)
            return target_path.exists()
        except Exception:
            return False
        finally:
            if presentation is not None:
                try:
                    presentation.Close()
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            if initialized:
                try:
                    pythoncom.CoUninitialize()
                except Exception:
                    pass

    def _infer_template_tags(self, name: str):
        source = str(name or "").lower()
        tags = ["模板", "ppt", "教学"]
        keyword_map = {
            "math": ["数学", "高数", "函数", "几何"],
            "function": ["函数", "图像", "推导"],
            "ratio": ["数学", "比例", "计算"],
            "angle": ["数学", "几何", "图形"],
            "geometric": ["几何", "结构", "理性"],
            "blackboard": ["课堂", "板书", "学术"],
            "cute": ["可爱", "卡通", "轻松"],
            "playful": ["互动", "活泼", "趣味"],
            "notebook": ["手账", "记录", "学习单"],
            "history": ["历史", "人文", "叙事"],
            "psychology": ["心理", "情绪", "案例"],
            "medical": ["医学", "科学", "知识"],
            "stem": ["科技", "综合", "探究"],
            "matisse": ["艺术", "创意", "插画"],
        }
        for keyword, mapped_tags in keyword_map.items():
            if keyword in source:
                tags.extend(mapped_tags)
        if "slidescarnival" in source:
            tags.extend(["演示", "教育"])
        if not any(tag in tags for tag in ["数学", "高数"]):
            tags.append("通用")
        return list(dict.fromkeys(tags))

    def _infer_template_scene(self, name: str):
        source = str(name or "").lower()
        if any(token in source for token in ["math", "function", "ratio", "angle", "geometric"]):
            return "数学/高数讲解、题型训练、概念建构"
        if any(token in source for token in ["psychology", "mental_health"]):
            return "心理与德育主题、讨论式课堂"
        if any(token in source for token in ["history", "medical", "stem"]):
            return "学科专题课、知识讲解与案例呈现"
        if any(token in source for token in ["cute", "playful", "doodles", "notebook"]):
            return "趣味课堂、互动导入、低龄或轻松表达"
        return "通用教学演示、课件设计"

    def handle_chat(self, session_id: str, session: dict, user_message: str):
        session["messages"].append({"role": "user", "content": user_message})
        self._rule_fill_slots(session, user_message)
        self._collect_creative_requests(session, user_message)

        fallback = self._fallback_clarification(session)
        llm_result = self.llm_client.chat_json(
            system_prompt=(
                "你是教师备课智能助理。"
                "请从用户消息中抽取 slots：course_theme、knowledge_points、key_difficulties、lesson_periods、style。"
                "允许从口语化表达、间接描述和上下文中推断字段。"
                "如果信息不全，只追问最关键的一个缺失字段。"
                "如果全部字段完整，则提示教师现在可以选择推荐模板，也可以继续补充课件/教案要求。"
                "严格返回 JSON：slots、assistant_reply、ready_to_generate。"
            ),
            user_prompt=json.dumps(
                {
                    "known_slots": session["slots"],
                    "user_message": user_message[:1200],
                    "required_fields": REQUIRED_FIELDS,
                    "creative_requests": session.get("creative_requests", []),
                },
                ensure_ascii=False,
            ),
            fallback=fallback,
        )

        for key, value in (llm_result.get("slots") or {}).items():
            if key in session["slots"] and str(value).strip():
                session["slots"][key] = str(value).strip()

        assistant_reply = str(llm_result.get("assistant_reply") or fallback["assistant_reply"]).strip()
        session["messages"].append({"role": "assistant", "content": assistant_reply})

        template_meta = self._template_meta(session)
        ready = self._all_slots_ready(session)
        if ready:
            session["template_picker_seen"] = True

        return {
            "assistant_reply": assistant_reply,
            "slots": session["slots"],
            "ready_to_generate": ready,
            "messages": session["messages"],
            "api_mode": "online" if self.llm_client.configured else "fallback",
            "template_suggestions": template_meta["suggestions"],
            "recommended_template": template_meta["recommended"],
            "selected_template": session.get("selected_template", ""),
            "template_picker_visible": ready,
            "template_catalog": self._template_catalog(template_meta["suggestions"]),
        }

    def generate_ppt(self, session_id: str, session: dict, selected_template: str = ""):
        self._update_progress(session, 5, "正在检索知识库...")
        package, retrievals, instruction_bundle = self._generate_package_core(
            session_id,
            session,
            revision="",
            selected_template=selected_template,
            include_interactive_media=True,
        )
        self._update_progress(session, 85, "正在导出 PPTX...")
        exported = self.export_service.build_ppt(package)
        self._update_progress(session, 95, "正在生成预览...")
        session["last_package"] = {"package": package, "files": exported}
        self._update_progress(session, 100, "生成完成")
        return self._build_package_response(session, package, retrievals, exported, instruction_bundle)

    def generate_doc(self, session_id: str, session: dict, selected_template: str = ""):
        self._update_progress(session, 5, "正在检索知识库...")
        package, retrievals, instruction_bundle = self._generate_package_core(
            session_id,
            session,
            revision="",
            selected_template=selected_template,
            include_interactive_media=False,
        )
        self._update_progress(session, 85, "正在导出 DOCX...")
        doc_safe_package = self._sanitize_package_for_doc_export(package)
        exported = self.export_service.build_doc(doc_safe_package)
        session["last_package"] = {"package": package, "files": exported}
        self._update_progress(session, 100, "生成完成")
        return self._build_package_response(session, package, retrievals, exported, instruction_bundle)

    def revise_teaching_package(self, session_id: str, session: dict, revision: str, selected_template: str = ""):
        self._collect_creative_requests(session, revision)
        self._update_progress(session, 5, "正在根据意见优化...")
        package, retrievals, instruction_bundle = self._generate_package_core(
            session_id,
            session,
            revision=revision,
            selected_template=selected_template,
            include_interactive_media=True,
        )
        package["applied_revision"] = revision
        self._update_progress(session, 80, "正在导出文件...")
        doc_safe_package = self._sanitize_package_for_doc_export(package)
        exported = {}
        exported.update(self.export_service.build_ppt(package))
        exported.update(self.export_service.build_doc(doc_safe_package))
        session["last_package"] = {"package": package, "files": exported}
        self._update_progress(session, 100, "优化完成")
        response = self._build_package_response(session, package, retrievals, exported, instruction_bundle)
        response["revision_applied"] = revision
        return response

    def _sanitize_package_for_doc_export(self, package: dict):
        safe_package = copy.deepcopy(package)
        safe_slides = []
        for slide in safe_package.get("slides") or []:
            slide_copy = dict(slide)
            bullets = []
            for bullet in slide_copy.get("bullets") or []:
                text, _ = self._normalize_bullet_item(bullet)
                if text:
                    bullets.append(text)
            slide_copy["bullets"] = bullets
            safe_slides.append(slide_copy)
        safe_package["slides"] = safe_slides
        return safe_package

    def _ensure_interactive_movie_embedding_hook(self):
        service = self.export_service

        if getattr(service, "_interactive_movie_hook_installed", False):
            return

        original_build_ppt = service._build_ppt
        original_add_interactive = service._add_interactive_slide
        original_add_content = service._add_content_slide

        def _wrapped_build_ppt(instance, file_path: Path, package: dict):
            instance._active_export_package = package
            return original_build_ppt(file_path, package)

        def _find_asset(instance, package: dict, slide_data: dict, slide_index: int | None = None):
            assets = package.get("interactive_assets") or []
            if not isinstance(assets, list) or not assets:
                return None
            if slide_index is not None:
                for item in assets:
                    if int(item.get("slide_index", -1)) == int(slide_index):
                        return item
            normalized_title = str(slide_data.get("title") or "").strip()
            if normalized_title:
                for item in assets:
                    if str(item.get("slide_title") or "").strip() == normalized_title:
                        return item
            return assets[0]

        def _poster_path(instance):
            from PIL import Image, ImageDraw

            poster_dir = instance.output_dir / "interactive_assets" / "posters"
            poster_dir.mkdir(parents=True, exist_ok=True)
            path = poster_dir / "default_interactive_poster.png"
            if path.exists():
                return path
            image = Image.new("RGB", (1280, 720), (42, 89, 182))
            draw = ImageDraw.Draw(image)
            draw.rectangle((60, 60, 1220, 660), outline=(255, 255, 255), width=4)
            draw.text((110, 180), "课堂知识演示视频", fill=(255, 255, 255))
            draw.text((110, 290), "播放方式：在 PPT 放映模式点击播放", fill=(235, 242, 255))
            image.save(path)
            return path

        def _style_paragraph(instance, paragraph, *, size_pt: float, bold: bool = False, color=None):
            styler = getattr(instance, "_apply_body_paragraph_style", None)
            if callable(styler):
                return styler(paragraph, size_pt=size_pt, bold=bold, color=color)
            paragraph.font.size = Pt(size_pt)
            paragraph.font.bold = bold
            paragraph.font.name = "Microsoft YaHei"
            if color is not None:
                paragraph.font.color.rgb = RGBColor(*color)

        def _add_open_demo_links(instance, slide):
            box = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(1.0),
                Inches(6.0),
                Inches(11.35),
                Inches(0.95),
            )
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*instance.theme["card"])
            box.line.color.rgb = RGBColor(*instance.theme["accent"])

            frame = box.text_frame
            frame.clear()

            title = frame.paragraphs[0]
            title.text = "延伸学习链接（点击可访问）"
            _style_paragraph(instance, title, size_pt=12, bold=True, color=instance.theme["accent"])

            for idx, (label, url) in enumerate(self.OPEN_DEMO_LINKS):
                p = frame.add_paragraph()
                p.text = ""
                run = p.add_run()
                run.text = f"{idx + 1}. {label}"
                run.hyperlink.address = url
                _style_paragraph(instance, p, size_pt=11, color=instance.theme["dark_text"])

        def _add_resource_links_slide(instance, prs, slide_data: dict):
            slide = instance._add_blank_slide(prs)
            instance._set_background(slide, instance.theme["soft"])
            instance._add_top_band(slide)
            instance._text_box(slide, 0.85, 0.55, 8.6, 0.8, slide_data.get("title") or "开源演示资源导航", 24, instance.theme["dark_text"], bold=True)

            resources = list(slide_data.get("resources") or [])[:4]
            if not resources:
                active_package = getattr(instance, "_active_export_package", {}) or {}
                resources = list(active_package.get("open_demo_references") or [])[:4]
            if not resources:
                resources = [
                    {
                        "name": site.get("name", ""),
                        "url": site.get("url", ""),
                        "summary": site.get("default_summary", "用于课堂演示与课后拓展。"),
                        "brief_intro": site.get("brief_intro", ""),
                        "teaching_hint": "；".join(site.get("suggestions", [])) or "建议结合本节主题进行定向检索。",
                    }
                    for site in self.OPEN_DEMO_SITES
                ]

            card_w = 5.52
            card_h = 2.50
            positions = [(0.95, 1.45), (6.83, 1.45), (0.95, 4.10), (6.83, 4.10)]
            for idx, item in enumerate(resources[:4]):
                left, top = positions[idx]
                card = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                    Inches(left),
                    Inches(top),
                    Inches(card_w),
                    Inches(card_h),
                )
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(*instance.theme["card"])
                card.line.color.rgb = RGBColor(*instance.theme["accent"])

                frame = card.text_frame
                frame.clear()

                name_p = frame.paragraphs[0]
                name_p.text = str(item.get("name") or f"资源 {idx + 1}")
                _style_paragraph(instance, name_p, size_pt=14, bold=True, color=instance.theme["accent"])

                intro_text = str(item.get("brief_intro") or "").strip()
                if intro_text:
                    intro_p = frame.add_paragraph()
                    intro_p.text = intro_text[:120]
                    _style_paragraph(instance, intro_p, size_pt=10, color=instance.theme["dark_text"])

                summary_p = frame.add_paragraph()
                summary_p.text = str(item.get("summary") or "该网站可用于数学可视化演示。")[:110]
                _style_paragraph(instance, summary_p, size_pt=11, color=instance.theme["dark_text"])

                hint_p = frame.add_paragraph()
                hint_p.text = f"课堂用法：{str(item.get('teaching_hint') or '作为拓展演示与讨论素材')}"[:120]
                _style_paragraph(instance, hint_p, size_pt=10, color=instance.theme["dark_text"])

                link_p = frame.add_paragraph()
                link_p.text = ""
                link_run = link_p.add_run()
                link_run.text = "🔗 打开链接"
                link_run.hyperlink.address = str(item.get("url") or "")
                _style_paragraph(instance, link_p, size_pt=11, bold=True, color=instance.theme["accent"])

        def _wrapped_add_interactive(instance, prs, slide_data: dict, package: dict | None = None, slide_index: int | None = None):
            resolved_package = package if isinstance(package, dict) else getattr(instance, "_active_export_package", {})
            asset = _find_asset(instance, resolved_package, slide_data, slide_index)
            if not asset:
                return original_add_interactive(prs, slide_data)

            video_path = Path(str(asset.get("video_path") or "").strip())
            if not video_path.exists() or video_path.suffix.lower() != ".mp4":
                return original_add_interactive(prs, slide_data)

            slide = instance._add_blank_slide(prs)
            instance._set_background(slide, instance.theme["soft"])
            instance._add_top_band(slide)
            instance._text_box(slide, 0.85, 0.55, 7.0, 0.8, slide_data.get("title") or "互动页", 24, instance.theme["dark_text"], bold=True)

            try:
                poster = _poster_path(instance)
                slide.shapes.add_movie(
                    str(video_path.resolve()),
                    Inches(1.0),
                    Inches(1.55),
                    Inches(11.35),
                    Inches(4.35),
                    str(poster.resolve()),
                    "video/mp4",
                )
                caption = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.0), Inches(5.95), Inches(11.35), Inches(0.38))
                caption.fill.solid()
                caption.fill.fore_color.rgb = RGBColor(*instance.theme["card"])
                caption.line.fill.background()
                frame = caption.text_frame
                frame.clear()
                p = frame.paragraphs[0]
                p.text = f"知识演示视频：{asset.get('slide_title') or slide_data.get('title') or '演示环节'}"
                _style_paragraph(instance, p, size_pt=11, color=instance.theme["dark_text"])
                _add_open_demo_links(instance, slide)
                return
            except Exception:
                return original_add_interactive(prs, slide_data)

        def _wrapped_add_content(instance, prs, slide_data: dict):
            if str(slide_data.get("layout") or "").strip().lower() == "resource_links":
                return _add_resource_links_slide(instance, prs, slide_data)
            return original_add_content(prs, slide_data)

        service._build_ppt = types.MethodType(_wrapped_build_ppt, service)
        service._add_interactive_slide = types.MethodType(_wrapped_add_interactive, service)
        service._add_content_slide = types.MethodType(_wrapped_add_content, service)
        service._interactive_movie_hook_installed = True

    def _generate_package_core(self, session_id: str, session: dict, revision: str, selected_template: str, include_interactive_media: bool = True):
        template_meta = self._template_meta(session)
        if selected_template and selected_template in self.template_names:
            session["selected_template"] = selected_template
        elif session.get("selected_template") and session.get("selected_template") not in self.template_names:
            session["selected_template"] = ""
        elif not session.get("selected_template"):
            session["selected_template"] = template_meta["recommended"]

        self._update_progress(session, 10, "正在检索知识库...")
        search_plan = self._compose_search_plan(session, revision)
        retrievals = self.rag_service.search(
            session_id=session_id,
            query=search_plan["query"],
            top_k=8,
            include_global=True,
            query_hints=search_plan["query_hints"],
        )
        self._update_progress(session, 20, "正在生成课件内容...")
        instruction_bundle = self._build_instruction_bundle(session, retrievals, revision, session["selected_template"])
        package = self._build_package(session=session, retrievals=retrievals, instruction_bundle=instruction_bundle)
        package = self._ensure_interactive_slide(package=package, instruction_bundle=instruction_bundle)
        has_open_resource = self._has_open_resource_intent(instruction_bundle)
        has_interactive = include_interactive_media and self._has_interactive_intent(instruction_bundle)

        if has_open_resource and has_interactive:
            self._update_progress(session, 45, "正在获取资源与生成演示...")
            with ThreadPoolExecutor(max_workers=2) as executor:
                demo_future = executor.submit(self._build_open_demo_references, instruction_bundle)
                asset_future = executor.submit(self._build_interactive_assets, session_id, package, instruction_bundle)
                open_demo_references = demo_future.result()
                interactive_assets = asset_future.result()
        elif has_open_resource:
            self._update_progress(session, 45, "正在获取开源演示资源...")
            open_demo_references = self._build_open_demo_references(instruction_bundle)
            interactive_assets = []
        elif has_interactive:
            self._update_progress(session, 45, "正在生成知识演示视频...")
            open_demo_references = []
            interactive_assets = self._build_interactive_assets(session_id, package, instruction_bundle)
        else:
            open_demo_references = []
            interactive_assets = []

        if open_demo_references:
            package["open_demo_references"] = open_demo_references
            package = self._ensure_open_demo_resource_slide(package=package, references=open_demo_references)
        if interactive_assets:
            package["interactive_assets"] = interactive_assets
        package["instruction_bundle"] = instruction_bundle
        return package, retrievals, instruction_bundle

    def _build_interactive_assets(self, session_id: str, package: dict, instruction_bundle: dict):
        interactive_slides = [
            {"index": idx, "slide": slide}
            for idx, slide in enumerate(package.get("slides") or [])
            if str(slide.get("layout") or "").strip().lower() == "interactive"
        ]
        if not interactive_slides or self.interactive_max_videos <= 0:
            return []
        if not self._has_interactive_intent(instruction_bundle):
            return []

        topics = self._interactive_topics(package=package, instruction_bundle=instruction_bundle)
        demo_type = self._interactive_demo_type(instruction_bundle)
        target_dir = self.interactive_assets_root / self._safe_token(session_id)
        target_dir.mkdir(parents=True, exist_ok=True)

        assets = []
        for idx, record in enumerate(interactive_slides[: self.interactive_max_videos]):
            slide_index = int(record["index"])
            slide = record["slide"]
            topic = topics[idx] if idx < len(topics) else topics[-1]
            token_seed = f"{session_id}:{slide_index}:{topic}:{demo_type}"
            token = self._safe_token(token_seed)

            try:
                html_path = Path(self.interactive_service.generate_demo(topic=topic, demo_type=demo_type, output_dir=target_dir))
                video_path = target_dir / f"interactive_{token}.mp4"
                rendered = self.video_renderer.render(
                    html_path=html_path,
                    output_path=video_path,
                    duration=self.interactive_clip_duration,
                )
                assets.append(
                    {
                        "slide_index": slide_index,
                        "slide_title": str(slide.get("title") or f"互动页 {slide_index + 1}"),
                        "topic": topic,
                        "demo_type": demo_type,
                        "game_type": demo_type,
                        "html_path": str(html_path.resolve()),
                        "video_path": str(Path(rendered).resolve()),
                    }
                )
            except Exception as exc:
                assets.append(
                    {
                        "slide_index": slide_index,
                        "slide_title": str(slide.get("title") or f"互动页 {slide_index + 1}"),
                        "topic": topic,
                        "demo_type": demo_type,
                        "game_type": demo_type,
                        "html_path": "",
                        "video_path": "",
                        "error": str(exc)[:180],
                    }
                )
        return assets

    def _ensure_interactive_slide(self, package: dict, instruction_bundle: dict):
        slides = list(package.get("slides") or [])
        has_interactive = any(str(item.get("layout") or "").strip().lower() == "interactive" for item in slides)
        if has_interactive or not self._has_interactive_intent(instruction_bundle):
            return package

        teacher_intent = instruction_bundle.get("teacher_intent") or {}
        topic = str(teacher_intent.get("course_theme") or "课堂知识点").strip() or "课堂知识点"
        creative_requests = [str(item).strip() for item in (teacher_intent.get("creative_requests") or []) if str(item).strip()]
        demo_type = self._interactive_demo_type(instruction_bundle)
        demo_map = {
            "concept": "概念渐进演示",
            "timeline": "过程推演演示",
            "graph": "参数动态图像演示",
        }
        demo_name = demo_map.get(demo_type, "知识演示")

        slides.append(
            {
                "title": "知识演示",
                "layout": "interactive",
                "bullets": [
                    {"text": f"演示类型：{demo_name}（围绕“{topic}”）", "section_hint": "main"},
                    {"text": "目标：通过动态可视化帮助学生快速理解核心概念与变化规律。", "section_hint": "side"},
                    {"text": f"教师要求：{'；'.join(creative_requests[:2]) or '强化概念理解与迁移应用'}", "section_hint": "note"},
                ],
            }
        )
        package["slides"] = slides
        return package

    def _ensure_open_demo_resource_slide(self, package: dict, references: list[dict]):
        slides = list(package.get("slides") or [])
        if not references:
            package["slides"] = slides
            return package

        if any(str(item.get("layout") or "").strip().lower() == "resource_links" for item in slides):
            package["slides"] = slides
            return package

        resource_slide = {
            "title": "开源演示资源导航",
            "layout": "resource_links",
            "resources": references[:4],
            "bullets": [
                {
                    "text": f"{item.get('name', '资源')}：{item.get('brief_intro') or item.get('summary', '')}",
                    "section_hint": "main",
                }
                for item in references[:4]
            ],
        }

        inserted = False
        new_slides = []
        for item in slides:
            if not inserted and str(item.get("layout") or "").strip().lower() == "summary":
                new_slides.append(resource_slide)
                inserted = True
            new_slides.append(item)
        if not inserted:
            new_slides.append(resource_slide)

        package["slides"] = new_slides
        return package

    def _build_open_demo_references(self, instruction_bundle: dict):
        teacher_intent = instruction_bundle.get("teacher_intent") or {}
        course_theme = str(teacher_intent.get("course_theme") or "").strip()
        knowledge = str(teacher_intent.get("knowledge_points") or "").strip()
        revision = str(teacher_intent.get("revision") or "").strip()
        creative = " ".join(teacher_intent.get("creative_requests") or [])
        focus_terms = self._extract_focus_terms(" ".join([revision, creative]))
        topic_parts = [item for item in [course_theme, knowledge] if item]
        if focus_terms:
            topic_parts.append(" ".join(focus_terms[:2]))
        topic = " ".join(topic_parts).strip() or "高等数学"
        topic = topic[:72]
        cache_key = self._safe_token(topic)
        if cache_key in self._site_reference_cache:
            return self._site_reference_cache[cache_key]

        refs = []
        site_results = {}
        with ThreadPoolExecutor(max_workers=min(len(self.OPEN_DEMO_SITES), 4)) as executor:
            future_to_idx = {
                executor.submit(self._build_single_site_ref, site, topic): idx
                for idx, site in enumerate(self.OPEN_DEMO_SITES)
            }
            for future in as_completed(future_to_idx):
                idx = future_to_idx[future]
                try:
                    site_results[idx] = future.result()
                except Exception:
                    site_results[idx] = {
                        "name": self.OPEN_DEMO_SITES[idx].get("name", ""),
                        "url": self.OPEN_DEMO_SITES[idx].get("url", ""),
                        "summary": self.OPEN_DEMO_SITES[idx].get("default_summary", "可用于课堂演示。"),
                        "brief_intro": self.OPEN_DEMO_SITES[idx].get("brief_intro", ""),
                        "teaching_hint": "建议按本节关键词检索并挑选一个演示用于课堂讲解。",
                    }
        for idx in sorted(site_results.keys()):
            refs.append(site_results[idx])

        self._site_reference_cache[cache_key] = refs
        return refs

    def _build_single_site_ref(self, site: dict, topic: str):
        target_url = self._find_precise_site_page_url(site=site, topic=topic)
        return {
            "name": site["name"],
            "url": target_url,
            "summary": self._site_summary(site, topic, target_url),
            "brief_intro": site.get("brief_intro", ""),
            "teaching_hint": self._site_teaching_hint(site, topic),
        }

    def _has_open_resource_intent(self, instruction_bundle: dict):
        teacher_intent = instruction_bundle.get("teacher_intent") or {}
        corpus = " ".join(
            [
                str(teacher_intent.get("knowledge_points") or ""),
                str(teacher_intent.get("revision") or ""),
                " ".join(teacher_intent.get("creative_requests") or []),
            ]
        ).lower()
        keywords = [
            "演示", "可视化", "演示视频", "知识演示", "链接", "网站", "资源",
            "geogebra", "desmos", "math3d", "wolfram",
        ]
        if any(token in corpus for token in keywords):
            return True
        return self._has_interactive_intent(instruction_bundle)

    def _build_site_url(self, site: dict, topic: str):
        template = str(site.get("search") or site.get("url") or "").strip()
        if not template:
            return ""
        q = self._build_site_query_text(topic=topic, site=site)
        q = re.sub(r"\s+", " ", q).strip()
        if "{query_plus}" in template:
            return template.format(query=quote(q), query_plus=quote_plus(q))
        if "{query}" in template:
            return template.format(query=quote(q))
        if "?" in template:
            return template
        if "geogebra.org" in template:
            return f"{template.rstrip('/')}/search/{quote(q)}"
        return template

    def _build_site_query_text(self, topic: str, site: dict):
        candidates = self._topic_query_candidates(topic=topic, site=site)
        if candidates:
            return candidates[0]
        return re.sub(r"\s+", " ", str(topic or "")).strip() or "calculus function"

    def _find_precise_site_page_url(self, site: dict, topic: str):
        domain = str(site.get("domain") or urlparse(str(site.get("url") or "")).netloc).lower().strip()
        fallback_url = self._build_site_url(site, topic)
        if not domain:
            return fallback_url

        curated_url = self._pick_curated_site_url(site=site, topic=topic)
        if curated_url:
            return curated_url

        queries = self._topic_query_candidates(topic=topic, site=site)
        templates = list(site.get("search_templates") or [])
        if not templates:
            templates = [str(site.get("search") or site.get("url") or "").strip()]

        candidates = {}
        seen = set()
        attempts = 0
        max_attempts = 5
        for template in templates[:3]:
            for q in queries[:3]:
                if attempts >= max_attempts:
                    break
                search_url = str(template or "").format(query=quote(q), query_plus=quote_plus(q))
                attempts += 1
                html = self._fetch_web_html(search_url)
                if not html:
                    continue
                for link in self._extract_links_from_html(search_url, html):
                    norm = link.strip()
                    if not norm or norm in seen:
                        continue
                    seen.add(norm)
                    if self._is_asset_url(norm):
                        continue
                    parsed = urlparse(norm)
                    netloc = parsed.netloc.lower()
                    if domain not in netloc:
                        continue
                    score = self._score_page_match(norm, topic, site_name=str(site.get("name") or ""), site=site)
                    if score <= 0:
                        continue
                    candidates[norm] = max(float(score), float(candidates.get(norm, 0)))

                # 对列表页进行二跳解析，尽量拿到站内详情页。
                seeds = sorted(candidates.items(), key=lambda item: item[1], reverse=True)[:3]
                for seed_url, seed_score in seeds:
                    if not self._is_listing_url(seed_url, site=site):
                        continue
                    deep_html = self._fetch_web_html(seed_url)
                    if not deep_html:
                        continue
                    for deep_link in self._extract_links_from_html(seed_url, deep_html):
                        deep_url = deep_link.strip()
                        if not deep_url or deep_url in seen:
                            continue
                        seen.add(deep_url)
                        if self._is_asset_url(deep_url):
                            continue
                        deep_parsed = urlparse(deep_url)
                        if domain not in deep_parsed.netloc.lower():
                            continue
                        deep_score = self._score_page_match(deep_url, topic, site_name=str(site.get("name") or ""), site=site)
                        if self._is_detail_url(deep_url, site=site):
                            deep_score += 4
                        if deep_score <= 0:
                            continue
                        candidates[deep_url] = max(float(deep_score), float(candidates.get(deep_url, 0)), float(seed_score))
            if attempts >= max_attempts:
                break

        if not candidates:
            return fallback_url

        ranked = sorted(candidates.items(), key=lambda item: item[1], reverse=True)
        enriched = []
        for url, score in ranked[:6]:
            title_text, _ = self._fetch_web_title_desc(url)
            title_bonus = self._score_title_match(title_text=title_text, topic=topic)
            detail_bonus = 2 if self._is_detail_url(url, site=site) else 0
            enriched.append((score + title_bonus + detail_bonus, url, title_text))
        if not enriched:
            return fallback_url

        enriched.sort(key=lambda item: item[0], reverse=True)
        if len(enriched) >= 2 and abs(enriched[0][0] - enriched[1][0]) <= 2.5:
            llm_pick = self._llm_pick_best_site_url(site=site, topic=topic, ranked_candidates=enriched[:4])
            if llm_pick:
                return llm_pick

        best = enriched[0][1]
        return best or fallback_url

    def _pick_curated_site_url(self, site: dict, topic: str):
        site_name = str(site.get("name") or "").lower()
        topic_text = str(topic or "").lower()
        focus_terms = set(self._extract_focus_terms(topic))

        def _has(*tokens):
            return any(token in topic_text or token in focus_terms for token in tokens)

        candidates = []
        if "geogebra" in site_name:
            if _has("3d", "geometry"):
                candidates = ["https://www.geogebra.org/3d", "https://www.geogebra.org/graphing"]
            else:
                candidates = ["https://www.geogebra.org/graphing", "https://www.geogebra.org/classic"]
        elif "desmos" in site_name:
            if _has("3d", "geometry", "surface"):
                candidates = ["https://www.desmos.com/3d", "https://www.desmos.com/calculator"]
            elif _has("ratio", "arithmetic"):
                candidates = ["https://www.desmos.com/fourfunction", "https://www.desmos.com/calculator"]
            else:
                candidates = ["https://www.desmos.com/calculator", "https://www.desmos.com/3d"]
        elif "math3d" in site_name:
            if _has("3d", "geometry", "surface"):
                candidates = ["https://www.math3d.org/gallery", "https://www.math3d.org/examples"]
            else:
                candidates = ["https://www.math3d.org/examples", "https://www.math3d.org/gallery"]
        elif "wolfram" in site_name:
            if _has("limit", "epsilon", "delta"):
                candidates = [
                    "https://demonstrations.wolfram.com/EpsilonDeltaDefinitionOfALimit/",
                    "https://demonstrations.wolfram.com/LimitOfAFunction/",
                    "https://demonstrations.wolfram.com/OneSidedLimits/",
                ]
            elif _has("derivative"):
                candidates = [
                    "https://demonstrations.wolfram.com/DerivativeAsTheLimitOfTheDifferenceQuotient/",
                    "https://demonstrations.wolfram.com/topics/Calculus.html",
                ]
            elif _has("integral"):
                candidates = [
                    "https://demonstrations.wolfram.com/CalculusAndAnalysis.html",
                    "https://demonstrations.wolfram.com/topics/Calculus.html",
                ]
            else:
                candidates = [
                    "https://demonstrations.wolfram.com/topics/Calculus.html",
                    "https://demonstrations.wolfram.com/CalculusAndAnalysis.html",
                ]

        for url in candidates:
            if self._url_is_reachable(url):
                return url
        return ""

    def _url_is_reachable(self, url: str):
        html = self._fetch_web_html(url)
        if len(html) < 200:
            return False
        lowered = html.lower()
        if "<html" not in lowered and "<!doctype" not in lowered:
            return False
        return True

    def _topic_query_candidates(self, topic: str, site: dict):
        base = re.sub(r"\s+", " ", str(topic or "")).strip() or "高等数学"
        site_hint = str(site.get("name") or "").lower()
        suggestions = list(site.get("suggestions") or [])
        mapped_terms = list(self._extract_focus_terms(base))
        llm_terms = self._llm_expand_topic_terms(base)
        short = " ".join(llm_terms[:3]).strip() or " ".join(mapped_terms[:3]).strip() or "calculus function"

        candidates = [
            f"{short} {site_hint}".strip(),
            short,
            "calculus function graph",
        ]
        for term in llm_terms[:3]:
            candidates.append(f"{term} {site_hint}".strip())
            candidates.append(term)
        if suggestions:
            candidates.append(" ".join(self._extract_focus_terms(str(suggestions[0]))[:3]) or short)
        deduped = []
        for item in candidates:
            val = re.sub(r"\s+", " ", item).strip()
            val = val[:48]
            if val and val not in deduped:
                deduped.append(val)
        return deduped

    def _llm_expand_topic_terms(self, topic: str):
        raw_topic = re.sub(r"\s+", " ", str(topic or "")).strip() or "高等数学"
        cache_key = self._safe_token(f"topic_terms::{raw_topic}")
        if cache_key in self._topic_term_cache:
            return list(self._topic_term_cache[cache_key])

        fallback_terms = self._extract_focus_terms(raw_topic)
        if not fallback_terms:
            fallback_terms = ["calculus", "function", "graph", "limit"]

        result = self.llm_client.chat_json(
            system_prompt=(
                "你是数学教学资源检索词助手。"
                "把中文教学主题转换为 4~8 个英文检索短词，"
                "仅输出 JSON：{\"keywords\": [\"...\"]}。"
                "关键词需适合 GeoGebra/Desmos/Math3D/Wolfram 检索。"
            ),
            user_prompt=json.dumps(
                {
                    "topic": raw_topic,
                    "fallback_terms": fallback_terms,
                    "constraints": ["english only", "2-4 words each", "math interactive demo oriented"],
                },
                ensure_ascii=False,
            ),
            fallback={"keywords": fallback_terms},
        )

        keywords = []
        for item in (result.get("keywords") or []):
            token = re.sub(r"[^a-z0-9\s\-]", " ", str(item).lower())
            token = re.sub(r"\s+", " ", token).strip()
            if not token:
                continue
            if token not in keywords:
                keywords.append(token)
        if not keywords:
            keywords = fallback_terms

        self._topic_term_cache[cache_key] = keywords[:8]
        return list(self._topic_term_cache[cache_key])

    def _extract_focus_terms(self, text: str):
        source = str(text or "")
        pairs = [
            ("极限", "limit"),
            ("导数", "derivative"),
            ("函数", "function"),
            ("图像", "graph"),
            ("参数", "parameter"),
            ("积分", "integral"),
            ("三维", "3d"),
            ("几何", "geometry"),
        ]
        terms = []
        lower = source.lower()
        for zh, en in pairs:
            if zh in source or en in lower:
                terms.append(en)
        return terms

    def _is_asset_url(self, url: str):
        lowered = str(url or "").lower()
        if any(token in lowered for token in ["/static/", "/assets/", "/img/", "/images/"]):
            return True
        if any(lowered.endswith(ext) for ext in [".png", ".jpg", ".jpeg", ".gif", ".svg", ".ico", ".css", ".js", ".woff", ".woff2", ".ttf"]):
            return True
        return False

    def _is_listing_url(self, url: str, site: dict):
        lowered = str(url or "").lower()
        listing_tokens = ["search", "query=", "/topics", "/browse", "/tag", "/tags", "/results", "/materials?", "format=rss"]
        if any(token in lowered for token in listing_tokens):
            return True
        for marker in (site.get("avoid_hints") or []):
            if str(marker).lower() in lowered:
                return True
        return False

    def _is_detail_url(self, url: str, site: dict):
        parsed = urlparse(str(url or ""))
        lowered = str(url or "").lower()
        if self._is_listing_url(lowered, site=site):
            return False
        for marker in (site.get("detail_hints") or []):
            if str(marker).lower() in lowered:
                return True
        path_segments = [seg for seg in parsed.path.split("/") if seg]
        if len(path_segments) >= 2:
            tail = path_segments[-1].lower()
            if any(ch.isdigit() for ch in tail) and len(tail) >= 4:
                return True
            if "-" in tail and len(tail) >= 8:
                return True
        return False

    def _fetch_web_html(self, url: str):
        target = str(url or "").strip()
        if not target.startswith("http"):
            return ""
        try:
            request = Request(target, headers={"User-Agent": "Mozilla/5.0 A04-Teaching-Agent"})
            with urlopen(request, timeout=2.2) as response:
                payload = response.read(300000)
            return payload.decode("utf-8", errors="ignore")
        except Exception:
            return ""

    def _extract_links_from_html(self, base_url: str, html: str):
        links = []
        for match in re.finditer(r"href=[\"']([^\"']+)[\"']", str(html or ""), flags=re.IGNORECASE):
            href = str(match.group(1) or "").strip()
            if not href or href.startswith("#"):
                continue
            if href.startswith("javascript:") or href.startswith("mailto:"):
                continue
            absolute = urljoin(base_url, href)
            absolute = self._unwrap_search_redirect(absolute)
            parsed = urlparse(absolute)
            if parsed.scheme not in {"http", "https"}:
                continue
            normalized = absolute.split("#", 1)[0]
            links.append(normalized)
        return links

    def _unwrap_search_redirect(self, url: str):
        target = str(url or "").strip()
        parsed = urlparse(target)
        query = parse_qs(parsed.query)

        for key in ["uddg", "url", "target", "dest", "destination", "r"]:
            val = (query.get(key) or [""])[0]
            normalized = unquote(str(val or "").strip())
            if normalized.startswith("http"):
                return normalized

        if "bing.com" in parsed.netloc.lower():
            raw_u = (query.get("u") or [""])[0]
            decoded = self._decode_bing_redirect(raw_u)
            if decoded:
                return decoded
        return target

    def _decode_bing_redirect(self, payload: str):
        raw = unquote(str(payload or "").strip())
        if not raw:
            return ""
        if raw.startswith("http"):
            return raw

        candidates = [raw]
        if raw.startswith("a1") and len(raw) > 2:
            candidates.append(raw[2:])

        for item in candidates:
            for padding in ["", "=", "==", "==="]:
                try:
                    decoded = base64.urlsafe_b64decode((item + padding).encode("utf-8")).decode("utf-8", errors="ignore").strip()
                except Exception:
                    continue
                if decoded.startswith("http"):
                    return decoded
        return ""

    def _score_page_match(self, url: str, topic: str, site_name: str, site: dict | None = None):
        score = 1.0
        url_lower = str(url).lower()
        if self._is_asset_url(url_lower):
            return -100

        if site and self._is_listing_url(url_lower, site=site):
            score -= 4
        if site and self._is_detail_url(url_lower, site=site):
            score += 5

        if url_lower.endswith("/") or url_lower.rstrip("/").count("/") <= 2:
            score -= 1.5
        if any(token in url_lower for token in ["uploads.", "cdn."]):
            score -= 3
        if any(token in url_lower for token in ["search", "login", "signup", "privacy", "tos"]):
            score -= 2

        if any(token in url_lower for token in ["/m/", "material", "activity", "lesson", "demonstration", "graph", "calculator", "surface", "3d"]):
            score += 4

        if "?" in url_lower:
            score += 1

        topic_tokens = [token for token in re.split(r"[，,、；;\s]+", str(topic or "").lower()) if len(token) >= 2]
        for token in topic_tokens[:8]:
            if token in url_lower:
                score += 4

        math_tokens = ["limit", "derivative", "function", "graph", "3d", "integral", "epsilon", "delta", "极限", "导数", "函数", "图像", "参数", "积分"]
        for token in math_tokens:
            if token in url_lower:
                score += 2

        if site_name and site_name.lower().replace(" ", "") in url_lower.replace("-", ""):
            score += 0.8

        return score

    def _score_title_match(self, title_text: str, topic: str):
        title = str(title_text or "").lower()
        if not title:
            return 0.0
        score = 0.0
        topic_terms = list(self._extract_focus_terms(topic)) + self._llm_expand_topic_terms(topic)[:5]
        seen = set()
        for term in topic_terms:
            for token in str(term).split():
                token = token.strip().lower()
                if len(token) < 3 or token in seen:
                    continue
                seen.add(token)
                if token in title:
                    score += 1.2
        return score

    def _llm_pick_best_site_url(self, site: dict, topic: str, ranked_candidates: list[tuple[float, str, str]]):
        if not self.llm_client.configured or not ranked_candidates:
            return ""

        candidate_payload = []
        for score, url, title in ranked_candidates[:4]:
            candidate_payload.append(
                {
                    "url": url,
                    "title": title,
                    "score": round(float(score), 3),
                }
            )

        fallback_url = candidate_payload[0]["url"]
        result = self.llm_client.chat_json(
            system_prompt=(
                "你是教学资源链接选择器。"
                "从候选链接中选择一个最适合课堂演示的详情页（不是搜索首页），"
                "返回 JSON：{\"best_url\": \"...\"}。"
            ),
            user_prompt=json.dumps(
                {
                    "site": site.get("name"),
                    "topic": topic,
                    "candidates": candidate_payload,
                    "rule": "prefer detail page over search/listing page",
                },
                ensure_ascii=False,
            ),
            fallback={"best_url": fallback_url},
        )
        selected = str(result.get("best_url") or "").strip()
        allowed = {item["url"] for item in candidate_payload}
        if selected in allowed:
            return selected
        return ""

    def _site_summary(self, site: dict, topic: str, target_url: str):
        default_text = str(site.get("default_summary") or "").strip() or "可用于课堂演示。"
        title_text, desc_text = self._fetch_web_title_desc(str(target_url or site.get("url") or ""))
        merged = "；".join([part for part in [title_text, desc_text] if part])
        if merged:
            merged = re.sub(r"\s+", " ", merged)
            return (merged[:72] + "…") if len(merged) > 72 else merged
        return default_text

    def _site_teaching_hint(self, site: dict, topic: str):
        topic_text = str(topic or "").lower()
        site_name = str(site.get("name") or "")
        if any(token in topic_text for token in ["函数", "图像", "导数", "极限", "参数"]):
            if "Desmos" in site_name:
                return "用于函数图像与参数滑块联动演示，强调变化趋势。"
            if "GeoGebra" in site_name:
                return "用于构造式可视化与课堂探究活动设计。"
            if "Math3D" in site_name:
                return "用于三维图形和曲面观察，连接平面与空间认知。"
        if "Wolfram" in site_name:
            return "用于补充拓展演示案例，安排课后探究任务。"
        suggestions = site.get("suggestions") or []
        if suggestions:
            return f"推荐方向：{suggestions[0]}"
        return "建议按本节关键词检索并挑选一个演示用于课堂讲解。"

    def _fetch_web_title_desc(self, url: str):
        cleaned_url = str(url or "").strip()
        if not cleaned_url.startswith("http"):
            return "", ""
        try:
            request = Request(cleaned_url, headers={"User-Agent": "Mozilla/5.0 A04-Teaching-Agent"})
            with urlopen(request, timeout=3.5) as response:
                payload = response.read(240000)
            text = payload.decode("utf-8", errors="ignore")
        except Exception:
            return "", ""

        title_match = re.search(r"<title[^>]*>(.*?)</title>", text, flags=re.IGNORECASE | re.DOTALL)
        desc_match = re.search(
            r"<meta[^>]+name=[\"']description[\"'][^>]+content=[\"'](.*?)[\"']",
            text,
            flags=re.IGNORECASE | re.DOTALL,
        )
        title = unescape(re.sub(r"\s+", " ", title_match.group(1))).strip() if title_match else ""
        desc = unescape(re.sub(r"\s+", " ", desc_match.group(1))).strip() if desc_match else ""
        return title[:44], desc[:88]

    def _has_interactive_intent(self, instruction_bundle: dict):
        teacher_intent = instruction_bundle.get("teacher_intent") or {}
        combined = " ".join(
            [
                str(teacher_intent.get("knowledge_points") or ""),
                str(teacher_intent.get("style") or ""),
                str(teacher_intent.get("revision") or ""),
                " ".join(teacher_intent.get("creative_requests") or []),
            ]
        )
        keywords = {"动画", "动效", "互动", "小游戏", "配对", "记忆", "闯关", "拖拽", "游戏", "演示", "可视化", "动态演示", "知识演示", "演示视频"}
        return any(token in combined for token in keywords)

    def _interactive_topics(self, package: dict, instruction_bundle: dict):
        teacher_intent = instruction_bundle.get("teacher_intent") or {}
        slots_theme = str(teacher_intent.get("course_theme") or "").strip()
        knowledge = str(teacher_intent.get("knowledge_points") or "").strip()
        creative_requests = [str(item).strip() for item in (teacher_intent.get("creative_requests") or []) if str(item).strip()]

        slide_topics = []
        for slide in package.get("slides") or []:
            if str(slide.get("layout") or "") != "interactive":
                continue
            title = str(slide.get("title") or "互动环节").strip()
            first_bullet = ""
            for bullet in slide.get("bullets") or []:
                text, _ = self._normalize_bullet_item(bullet)
                if text:
                    first_bullet = text
                    break
            composed = " / ".join([part for part in [slots_theme, knowledge, title, first_bullet] if part])
            slide_topics.append(composed or "课堂知识演示")

        if slide_topics:
            return slide_topics
        if creative_requests:
            return creative_requests[:1]
        base = " / ".join([part for part in [slots_theme, knowledge] if part]).strip()
        return [base or "课堂知识演示"]

    def _interactive_demo_type(self, instruction_bundle: dict):
        teacher_intent = instruction_bundle.get("teacher_intent") or {}
        combined = " ".join(
            [
                str(teacher_intent.get("revision") or ""),
                " ".join(teacher_intent.get("creative_requests") or []),
            ]
        )
        if any(token in combined for token in ["步骤", "过程", "推导", "时间线", "流程"]):
            return "timeline"
        if any(token in combined for token in ["函数", "图像", "曲线", "参数", "坐标", "变化"]):
            return "graph"
        return "concept"

    def _safe_token(self, raw: str):
        return hashlib.md5(str(raw).encode("utf-8")).hexdigest()[:14]

    def _compose_search_plan(self, session: dict, revision: str):
        slots = session["slots"]
        query_parts = [slots.get("course_theme", ""), slots.get("knowledge_points", ""), slots.get("key_difficulties", "")]
        if revision:
            query_parts.append(revision)
        query_parts.extend(session.get("creative_requests", [])[:3])
        query_parts = [item.strip() for item in query_parts if str(item).strip()]

        keywords = []
        desired_chunk_types = ["knowledge", "overview"]
        if slots.get("knowledge_points"):
            keywords.extend(re.split(r"[，、；;\s]+", slots["knowledge_points"]))
        if slots.get("key_difficulties"):
            keywords.extend(re.split(r"[，、；;\s]+", slots["key_difficulties"]))
        if any(token in f"{revision}{' '.join(session.get('creative_requests', []))}" for token in ["案例", "情境", "应用", "例题"]):
            desired_chunk_types.append("case")
            keywords.extend(["案例", "情境", "应用"])
        if any(token in f"{revision}{slots.get('style', '')}" for token in ["风格", "版式", "排版", "模板"]):
            desired_chunk_types.append("style")
            keywords.extend(["风格", "版式", "排版"])
        if any(token in f"{revision}{' '.join(session.get('creative_requests', []))}" for token in ["动画", "互动", "游戏", "闯关"]):
            keywords.extend(["互动", "动画", "游戏"])

        return {
            "query": " ".join(query_parts),
            "query_hints": {
                "keywords": [token for token in keywords if token][:8],
                "desired_chunk_types": list(dict.fromkeys(desired_chunk_types)),
            },
        }

    def _build_instruction_bundle(self, session: dict, retrievals: list[dict], revision: str, template_choice: str):
        session_documents = []
        for doc in session.get("documents", [])[-8:]:
            session_documents.append(
                {
                    "name": doc.get("original_name", "未命名资料"),
                    "summary": (doc.get("summary") or "")[:150],
                    "knowledge_structure": doc.get("knowledge_structure", [])[:3],
                    "cases": doc.get("cases", [])[:2],
                    "content_style": doc.get("content_style", ""),
                }
            )

        local_kb_hits = []
        uploaded_hits = []
        for item in retrievals:
            target = local_kb_hits if item.get("scope") == "global" else uploaded_hits
            target.append(
                {
                    "source": item.get("source_name", ""),
                    "title": item.get("title", ""),
                    "chunk_type": item.get("chunk_type", ""),
                    "score": item.get("score", 0),
                    "summary": item.get("summary", ""),
                    "excerpt": item.get("text", "")[:180],
                    "use": self._usage_hint(item),
                }
            )

        directives = [
            "以教师意图为主线组织内容，不要只堆砌资料片段。",
            "优先吸收上传资料中的知识结构、案例、风格信号，再用本地专业知识库补足专业讲解深度。",
            "所有页面文案要适合投屏阅读，每页控制在 3-5 条要点内。",
            "若教师提出动画创意或互动小游戏，必须落到具体可执行的课堂设计。",
        ]
        if revision:
            directives.append(f"必须落实教师本轮修改意见：{revision[:120]}")

        return {
            "teacher_intent": {
                "course_theme": session["slots"].get("course_theme", ""),
                "knowledge_points": session["slots"].get("knowledge_points", ""),
                "key_difficulties": session["slots"].get("key_difficulties", ""),
                "lesson_periods": session["slots"].get("lesson_periods", ""),
                "style": session["slots"].get("style", ""),
                "selected_template": template_choice,
                "creative_requests": session.get("creative_requests", [])[:4],
                "revision": revision[:200],
            },
            "teacher_context_digest": self._conversation_digest(session),
            "uploaded_materials": session_documents,
            "uploaded_material_hits": uploaded_hits[:4],
            "local_knowledge_hits": local_kb_hits[:4],
            "generation_directives": directives,
            "output_requirements": {
                "ppt_structure": ["封面", "目录", "正文", "案例或练习", "小结"],
                "readability": "高可读、短句、适合课堂投屏",
                "template_scene": self.template_scenes.get(template_choice, ""),
            },
        }

    def _usage_hint(self, item: dict):
        chunk_type = item.get("chunk_type", "")
        if chunk_type == "knowledge":
            return "用于补充概念定义、原理推导或知识结构"
        if chunk_type == "case":
            return "用于补充例题、情境案例或课堂练习"
        if chunk_type == "style":
            return "用于吸收表达顺序、排版风格或呈现结构"
        return "用于补充课堂背景信息和表达线索"

    def _build_package(self, session: dict, retrievals: list[dict], instruction_bundle: dict):
        slots = session["slots"]
        references = self._build_reference_notes(instruction_bundle)
        creative_requests = session.get("creative_requests", [])
        template_constraints = {
            "max_boxes_per_slide": 7,
            "min_boxes_per_slide": 5,
            "main_box_max_chars": 35,
            "side_box_max_chars": 25,
            "note_box_max_chars": 15,
            "total_chars_per_slide": [120, 180],
        }
        fallback = self._fallback_package(
            slots=slots,
            references=references,
            revision=instruction_bundle["teacher_intent"].get("revision", ""),
            creative_requests=creative_requests,
            template_choice=instruction_bundle["teacher_intent"]["selected_template"],
        )

        llm_result = self.llm_client.chat_json(
            system_prompt=(
                "你是资深教研员和课件设计师。"
                "请把教师意图、上传资料、本地专业知识库命中内容进行有效融合，形成完整课件。"
                "优先保持教学逻辑完整：导入-建构-案例/练习-互动-小结。"
                "不得堆砌检索原文，不得直接照抄资料。"
                "【重要：模板适配约束】"
                "请严格遵循模板文本框容量与分区语义。"
                "默认每页按 5-7 个文本框组织内容，并区分主内容区、侧边栏、备注区。"
                "主内容区（main）每条不超过35字；侧边栏（side）每条不超过25字；备注区（note）每条不超过15字。"
                "每页总字数（含标题）控制在120-180字。"
                "为保证课堂可讲性，agenda/content/case/summary 页默认输出 5-7 条要点，interactive 页输出 3-4 条（含规则/题型/反馈）。"
                "每个核心知识点尽量给出：定义或结论 + 方法步骤 + 易错提醒/应用提示。"
                "【内容结构要求】"
                "slides[].bullets[] 必须是对象，格式为 {text, section_hint}。"
                "section_hint 仅允许 main、side、note。"
                "同一概念的多句内容必须尽量放在同一个 section_hint 中，避免拆散。"
                "严格返回 JSON：title、theme_style、slides、closing_points、references。"
                "slides 每项必须包含 title、layout、bullets，layout 仅可为 cover、agenda、content、case、summary、interactive。"
            ),
            user_prompt=json.dumps({"instruction_bundle": instruction_bundle, "template_constraints": template_constraints}, ensure_ascii=False),
            fallback=fallback,
        )

        package = dict(fallback)
        package.update(
            {
                "title": llm_result.get("title") or fallback["title"],
                "theme_style": llm_result.get("theme_style") or fallback["theme_style"],
                "slides": self._enrich_slides_density(
                    self._normalize_slides_with_section_hints(llm_result.get("slides") or fallback["slides"]),
                    slots=slots,
                ),
                "closing_points": llm_result.get("closing_points") or fallback["closing_points"],
                "references": llm_result.get("references") or fallback["references"],
            }
        )
        package["slides"] = self._prepare_slides_for_export_compat(package.get("slides") or [])
        package["summary"] = slots.copy()
        package["selected_template"] = instruction_bundle["teacher_intent"]["selected_template"]
        package["template_constraints"] = template_constraints
        return package

    def _normalize_slides_with_section_hints(self, slides: list[dict] | None):
        normalized_slides = []
        allowed_layouts = {"cover", "agenda", "content", "case", "summary", "interactive"}

        for slide in slides or []:
            layout = str(slide.get("layout") or "content").strip().lower()
            if layout not in allowed_layouts:
                layout = "content"
            normalized_slide = {
                "title": str(slide.get("title") or "未命名页面").strip() or "未命名页面",
                "layout": layout,
                "bullets": [],
            }
            for bullet in slide.get("bullets") or []:
                text, hint = self._normalize_bullet_item(bullet)
                if not text:
                    continue
                normalized_slide["bullets"].append({"text": text, "section_hint": hint})
            normalized_slides.append(normalized_slide)

        return normalized_slides

    def _enrich_slides_density(self, slides: list[dict], slots: dict):
        min_items = {
            "cover": 3,
            "agenda": 5,
            "content": 6,
            "case": 6,
            "interactive": 4,
            "summary": 6,
        }

        knowledge = str(slots.get("knowledge_points") or "").strip()
        difficulties = str(slots.get("key_difficulties") or "").strip()

        def _append_unique(buffer: list[dict], text: str, hint: str = "main"):
            clean_text = re.sub(r"\s+", " ", str(text or "")).strip()
            if not clean_text:
                return
            if any((item.get("text") or "") == clean_text for item in buffer):
                return
            buffer.append({"text": clean_text, "section_hint": hint if hint in {"main", "side", "note"} else "main"})

        enriched_slides = []
        for slide in slides or []:
            layout = str(slide.get("layout") or "content")
            target = min_items.get(layout, 5)
            current = list(slide.get("bullets") or [])

            # 第一层：拆分长句，增加可讲述颗粒度。
            if len(current) < target:
                for item in list(current):
                    text = str(item.get("text") or "")
                    hint = str(item.get("section_hint") or "main")
                    segments = [re.sub(r"\s+", " ", seg).strip() for seg in re.split(r"[；;。]", text) if re.sub(r"\s+", " ", seg).strip()]
                    for seg in segments:
                        if len(seg) >= 6:
                            _append_unique(current, seg, hint)
                        if len(current) >= target:
                            break
                    if len(current) >= target:
                        break

            # 第二层：基于课程槽位补充“方法/误区/反馈”信息。
            if len(current) < target:
                layout_pool = {
                    "agenda": [
                        "目标：本节完成概念理解、方法掌握与应用表达",
                        "结构：导入→概念建构→例题演练→互动反馈→小结",
                        "检查：每部分结束后设置1个快速判断题",
                    ],
                    "content": [
                        f"知识主线：{knowledge}" if knowledge else "知识主线：概念—方法—应用三段推进",
                        f"难点提醒：{difficulties}" if difficulties else "难点提醒：建模思路与符号规范需同步把握",
                        "方法步骤：先读条件，再建模，再计算并检验",
                        "易错提醒：关注符号方向、区间端点和单位一致",
                    ],
                    "case": [
                        "案例结构：问题→思路→解答→反思",
                        "思路提示：先确定已知量、未知量与约束关系",
                        "迁移提问：条件改变时结果如何调整",
                    ],
                    "interactive": [
                        "规则：小组协作答题，按正确率与表达打分",
                        "题型：概念判断 + 计算题 + 应用解释",
                        "反馈：教师即时点评并归纳共性错误",
                    ],
                    "summary": [
                        "一句话总结：从定义出发，落实到可解释的应用",
                        "方法回顾：审题→建模→求解→验证",
                        "下节衔接：在综合场景中提升迁移能力",
                    ],
                }
                for candidate in layout_pool.get(layout, layout_pool["content"]):
                    _append_unique(current, candidate, "main")
                    if len(current) >= target:
                        break

            slide_copy = {
                "title": slide.get("title") or "未命名页面",
                "layout": layout,
                "bullets": current[:8],
            }
            enriched_slides.append(slide_copy)

        return enriched_slides

    def _normalize_bullet_item(self, bullet) -> tuple[str, str]:
        allowed_hints = {"main", "side", "note"}

        def _clean(v):
            return re.sub(r"\s+", " ", str(v or "")).strip()

        text = ""
        hint = "main"

        if isinstance(bullet, dict):
            text = _clean(bullet.get("text"))
            hint = _clean(bullet.get("section_hint") or "main").lower()
        else:
            raw = _clean(bullet)
            parsed = None
            # 兼容 LLM 返回 "{'text': '...', 'section_hint': 'main'}" 这种字典字符串。
            if raw.startswith("{") and raw.endswith("}"):
                try:
                    parsed = ast.literal_eval(raw)
                except Exception:
                    parsed = None
                if isinstance(parsed, dict):
                    text = _clean(parsed.get("text"))
                    hint = _clean(parsed.get("section_hint") or "main").lower()
                else:
                    text = raw
            else:
                text = raw

        if hint not in allowed_hints:
            hint = "main"
        return text, hint

    def _prepare_slides_for_export_compat(self, slides: list[dict]):
        # 兼容旧版 DocumentExportService（仅支持 bullets 为字符串列表）。
        if hasattr(self.export_service, "_bullet_text"):
            return slides

        flattened = []
        for slide in slides or []:
            flat_slide = {
                "title": slide.get("title") or "未命名页面",
                "layout": slide.get("layout") or "content",
                "bullets": [],
            }
            section_hints = []
            for bullet in slide.get("bullets") or []:
                text, hint = self._normalize_bullet_item(bullet)
                if not text:
                    continue
                flat_slide["bullets"].append(text)
                section_hints.append(hint)
            if section_hints:
                flat_slide["section_hints"] = section_hints
            flattened.append(flat_slide)
        return flattened

    def _build_package_response(self, session: dict, package: dict, retrievals: list[dict], downloads: dict, instruction_bundle: dict):
        template_meta = self._template_meta(session)
        doc_safe_package = self._sanitize_package_for_doc_export(package)
        slide_thumbnail_urls = []
        if downloads.get("pptx"):
            slide_thumbnail_urls = self.export_service.build_ppt_slide_thumbnails(downloads["pptx"])
        return {
            "package": package,
            "downloads": downloads,
            "ppt_preview": self._build_ppt_preview(package, slide_thumbnail_urls),
            "doc_preview": self.export_service.build_doc_preview(doc_safe_package),
            "api_mode": "online" if self.llm_client.configured else "fallback",
            "template_suggestions": template_meta["suggestions"],
            "recommended_template": template_meta["recommended"],
            "selected_template": session.get("selected_template", ""),
            "reference_digest": self._build_reference_digest(retrievals),
            "instruction_bundle_summary": {
                "uploaded_material_count": len(instruction_bundle.get("uploaded_materials", [])),
                "uploaded_hit_count": len(instruction_bundle.get("uploaded_material_hits", [])),
                "local_kb_hit_count": len(instruction_bundle.get("local_knowledge_hits", [])),
                "directives": instruction_bundle.get("generation_directives", [])[:4],
            },
            "template_catalog": self._template_catalog(template_meta["suggestions"]),
        }

    def _build_ppt_preview(self, package: dict, slide_thumbnail_urls: list[str] | None = None):
        thumb_urls = slide_thumbnail_urls or []
        slides = []
        for index, slide in enumerate(package.get("slides", [])[:12]):
            bullets = []
            for bullet in (slide.get("bullets") or [])[:4]:
                text, _ = self._normalize_bullet_item(bullet)
                bullets.append(text[:78] + ("…" if len(text) > 78 else ""))
            slides.append(
                {
                    "title": str(slide.get("title") or "未命名页面")[:40],
                    "layout": str(slide.get("layout") or "content"),
                    "bullets": bullets,
                    "thumbnail_url": thumb_urls[index] if index < len(thumb_urls) else "",
                }
            )
        return {"slides": slides}

    def build_session_snapshot(self, session_id: str, session: dict):
        template_meta = self._template_meta(session)
        snapshot = {
            "session_id": session_id,
            "messages": session.get("messages", []),
            "slots": session.get("slots", {}),
            "documents": session.get("documents", []),
            "selected_template": session.get("selected_template", ""),
            "template_suggestions": template_meta["suggestions"],
            "recommended_template": template_meta["recommended"],
            "template_catalog": self._template_catalog(template_meta["suggestions"]),
            "has_result": bool(session.get("last_package")),
        }
        if not session.get("last_package"):
            return snapshot

        package = (session.get("last_package") or {}).get("package") or {}
        doc_safe_package = self._sanitize_package_for_doc_export(package)
        downloads = (session.get("last_package") or {}).get("files") or {}
        slide_thumbnail_urls = []
        if downloads.get("pptx"):
            slide_thumbnail_urls = self.export_service.build_ppt_slide_thumbnails(downloads["pptx"])
        snapshot["result"] = {
            "package": package,
            "downloads": downloads,
            "ppt_preview": self._build_ppt_preview(package, slide_thumbnail_urls),
            "doc_preview": self.export_service.build_doc_preview(doc_safe_package),
            "api_mode": "online" if self.llm_client.configured else "fallback",
            "template_suggestions": template_meta["suggestions"],
            "recommended_template": template_meta["recommended"],
            "selected_template": session.get("selected_template", ""),
            "template_catalog": self._template_catalog(template_meta["suggestions"]),
            "reference_digest": [
                {
                    "title": "资料融入摘要",
                    "score": "",
                    "summary": item,
                    "highlights": ["历史会话恢复"],
                }
                for item in (package.get("references") or [])[:6]
            ],
        }
        return snapshot

    def _template_meta(self, session: dict):
        style = str(session["slots"].get("style", ""))
        creative = " ".join(session.get("creative_requests", []))
        corpus = f"{style} {creative} {session['slots'].get('course_theme', '')} {session['slots'].get('knowledge_points', '')}"
        scores = []
        for profile in self.template_profiles:
            score = 0
            if profile["name"] in style:
                score += 10
            score += sum(2 for tag in profile["tags"] if tag in corpus)
            if profile.get("source") == "external":
                score += 2
                if any(tag in (" ".join(profile.get("tags", []))) for tag in ["数学", "高数", "几何", "函数"]):
                    score += 2
            if any(keyword in creative for keyword in ["动画", "小游戏", "互动"]) and profile["name"] in {"卡通风", "可爱风", "插画风", "未来感"}:
                score += 3
            scores.append((score, profile["name"]))
        index_map = {name: idx for idx, name in enumerate(self.template_names)}
        scores.sort(key=lambda item: (-item[0], index_map.get(item[1], 9999)))
        suggestions = [name for _, name in scores[:6]]
        recommended = suggestions[0] if suggestions else "教育蓝"
        return {"suggestions": suggestions, "recommended": recommended}

    def _template_catalog(self, priority_templates: list[str]):
        ordered = []
        priority = list(dict.fromkeys(priority_templates))
        remaining = [name for name in self.template_names if name not in priority]
        for name in priority + remaining:
            profile = next((item for item in self.template_profiles if item["name"] == name), {})
            ordered.append(
                {
                    "name": name,
                    "scene": self.template_scenes.get(name, ""),
                    "recommended": name in priority[:3],
                    "thumbnail_url": profile.get("thumbnail_url", ""),
                }
            )
        return ordered

    def _build_reference_notes(self, instruction_bundle: dict):
        notes = []
        for item in instruction_bundle.get("uploaded_materials", [])[:6]:
            note = [item["name"]]
            if item.get("summary"):
                note.append(f"摘要：{item['summary']}")
            if item.get("knowledge_structure"):
                note.append(f"知识结构：{'；'.join(item['knowledge_structure'])}")
            if item.get("cases"):
                note.append(f"案例：{'；'.join(item['cases'])}")
            if item.get("content_style"):
                note.append(f"风格：{item['content_style']}")
            notes.append(" | ".join(note)[:280])

        for group_name in ("uploaded_material_hits", "local_knowledge_hits"):
            for item in instruction_bundle.get(group_name, [])[:4]:
                notes.append(f"{item['source']} | {item['use']} | {item['excerpt'][:120]}")
        return notes or ["暂无参考资料，按教师需求直接生成。"]

    def _build_reference_digest(self, retrievals: list[dict]):
        digest = []
        seen = set()
        for item in retrievals:
            key = (item.get("scope"), item.get("source_name"), item.get("title"))
            if key in seen:
                continue
            seen.add(key)
            title_prefix = "专业知识库" if item.get("scope") == "global" else "上传资料"
            digest.append(
                {
                    "title": f"{title_prefix} · {item.get('source_name', '')}",
                    "score": item.get("score", ""),
                    "summary": item.get("summary", "") or item.get("text", ""),
                    "highlights": [self._usage_hint(item), f"片段标题：{item.get('title', '正文')}"],
                }
            )
        return digest

    def _conversation_digest(self, session: dict):
        digest = []
        for item in session["messages"][-6:]:
            role = "教师" if item["role"] == "user" else "助理"
            digest.append(f"{role}：{item['content'][:120]}")
        return digest

    def _fallback_clarification(self, session: dict):
        missing = self._missing_fields(session)
        if missing:
            assistant_reply = f"我已经记录了部分教学需求。为了继续完善，请再补充“{REQUIRED_FIELDS[missing[0]]}”。"
        else:
            recommended = self._template_meta(session)["recommended"]
            assistant_reply = (
                f"教学需求已经完整。现在可以选择推荐模板“{recommended}”，"
                "也可以继续补充希望加入的动画创意、互动小游戏或教案要求。"
            )
        return {"slots": session["slots"], "assistant_reply": assistant_reply, "ready_to_generate": not missing}

    def _fallback_package(self, slots: dict, references: list[str], revision: str, creative_requests: list[str], template_choice: str):
        theme = slots.get("course_theme") or "课堂主题"
        points = slots.get("knowledge_points") or "核心知识点"
        difficulties = slots.get("key_difficulties") or "重难点"
        periods = slots.get("lesson_periods") or "1课时"
        style = template_choice or slots.get("style") or "教育蓝"
        revision_note = revision or "无"

        slides = [
            {
                "title": theme,
                "layout": "cover",
                "bullets": [
                    {"text": f"知识点：{points}", "section_hint": "main"},
                    {"text": f"课时安排：{periods}", "section_hint": "side"},
                    {"text": f"版式模板：{style}", "section_hint": "note"},
                ],
            },
            {
                "title": "学习导航",
                "layout": "agenda",
                "bullets": [
                    {"text": "情境导入", "section_hint": "main"},
                    {"text": "概念建构", "section_hint": "main"},
                    {"text": "例题讲解", "section_hint": "side"},
                    {"text": "互动练习", "section_hint": "side"},
                    {"text": "课堂小结", "section_hint": "note"},
                ],
            },
            {
                "title": "情境导入",
                "layout": "content",
                "bullets": [
                    {"text": "用真实问题引出本课主题。", "section_hint": "main"},
                    {"text": f"唤醒与“{points}”相关的前置知识。", "section_hint": "main"},
                    {"text": "用一个核心问题串起整节课。", "section_hint": "side"},
                ],
            },
            {
                "title": "知识结构梳理",
                "layout": "content",
                "bullets": [
                    {"text": f"围绕“{points}”拆分知识层级。", "section_hint": "main"},
                    {"text": f"重点突破：{difficulties}。", "section_hint": "note"},
                    {"text": "用板书结构或流程框帮助学生建立联系。", "section_hint": "side"},
                ],
            },
            {
                "title": "案例与练习",
                "layout": "case",
                "bullets": [
                    {"text": "引入资料中的典型案例。", "section_hint": "main"},
                    {"text": "设计基础到提升的练习链。", "section_hint": "side"},
                    {"text": "提醒学生说明方法选择与易错点。", "section_hint": "note"},
                ],
            },
        ]
        if creative_requests:
            slides.append(
                {
                    "title": "互动创意设计",
                    "layout": "interactive",
                    "bullets": [
                        {"text": "知识点动画创意：用逐步显现、遮罩揭示或移动路径突出概念变化。", "section_hint": "main"},
                        {"text": "互动小游戏：设计配对、闯关、拖拽排序或抢答环节，服务于知识点掌握。", "section_hint": "side"},
                        {"text": f"教师附加要求：{'；'.join(creative_requests[:3])}", "section_hint": "note"},
                    ],
                }
            )
        slides.append(
            {
                "title": "课堂小结",
                "layout": "summary",
                "bullets": [
                    {"text": "回顾本节课的核心概念、方法与误区。", "section_hint": "main"},
                    {"text": "用出口任务检查达成度。", "section_hint": "side"},
                    {"text": f"本次修改要求：{revision_note}", "section_hint": "note"},
                ],
            }
        )
        return {
            "title": f"{theme} 互动式教学课件",
            "theme_style": style,
            "slides": slides,
            "closing_points": [f"本课围绕“{theme}”建立清晰知识链路。", f"重点突破“{difficulties}”。", "建议导出后结合班级学情继续微调活动页。"],
            "references": references[:8],
        }

    def _rule_fill_slots(self, session: dict, message: str):
        text = str(message or "")
        patterns = {
            "course_theme": [r"(?:课程主题|主题|课题|课程)[：: ]?([^\n；;，,]+)"],
            "knowledge_points": [r"(?:知识点|核心知识|教学内容)[：: ]?([^\n；;]+)"],
            "key_difficulties": [r"(?:重难点|重点|难点)[：: ]?([^\n；;]+)"],
            "lesson_periods": [r"(?:课时安排|课时)[：: ]?([^\n；;]+)", r"(\d+\s*课时)", r"(\d+\s*分钟)"],
            "style": [r"(?:课件风格|风格|模板)[：: ]?([^\n；;]+)"],
        }
        for slot, regexes in patterns.items():
            for pattern in regexes:
                match = re.search(pattern, text)
                if match:
                    session["slots"][slot] = match.group(1).strip("。；;，, ")
                    break

        if not session["slots"].get("course_theme", "").strip():
            match = re.search(r"(?:这节课讲|本节课讲|准备讲|主要讲|我们讲)\s*([^\n，。；;]{2,30})", text)
            if match:
                session["slots"]["course_theme"] = match.group(1).strip("。；;，, ")

        if not session["slots"].get("knowledge_points", "").strip():
            match = re.search(r"(?:重点放在|核心是|主要是|会讲到|包括)\s*([^\n。；;]{3,60})", text)
            if match:
                session["slots"]["knowledge_points"] = match.group(1).strip("。；;，, ")

        if not session["slots"].get("key_difficulties", "").strip():
            match = re.search(r"(?:难点在于|学生容易|学生常错|容易混淆|最难的是)\s*([^\n。；;]{3,60})", text)
            if match:
                session["slots"]["key_difficulties"] = match.group(1).strip("。；;，, ")

        if not session["slots"].get("lesson_periods", "").strip():
            match = re.search(r"(\d+\s*(?:课时|分钟|min|节课))", text, re.IGNORECASE)
            if match:
                session["slots"]["lesson_periods"] = match.group(1).strip("。；;，, ")

        if not session["slots"].get("style", "").strip():
            match = re.search(r"(?:想要|希望|偏好|做成)\s*([^\n。；;]{2,20})(?:风|风格|模板)", text)
            if match:
                session["slots"]["style"] = f"{match.group(1).strip('。；;，, ')}风"

    def _collect_creative_requests(self, session: dict, text: str):
        compact = str(text or "").strip()
        if compact and any(keyword in compact for keyword in self.CREATIVE_KEYWORDS):
            if compact not in session["creative_requests"]:
                session["creative_requests"].append(compact[:200])

    def _missing_fields(self, session: dict):
        return [key for key, value in session["slots"].items() if not str(value).strip()]

    def _all_slots_ready(self, session: dict):
        return not self._missing_fields(session)
