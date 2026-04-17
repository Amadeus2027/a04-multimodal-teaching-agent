from __future__ import annotations

from pathlib import Path
import io
import logging
import re
import uuid

import easyocr
import numpy as np
from docx import Document as DocxDocument
from PIL import Image
from pptx import Presentation
from PyPDF2 import PdfReader
from werkzeug.utils import secure_filename

from services.audio_transcriber import OfflineAudioTranscriber

try:
    import cv2
except Exception:  # pragma: no cover
    cv2 = None


class FileParserService:
    """Parses PDF, image, Word, and video files into compact teaching notes."""

    ALLOWED_SUFFIXES = {
        ".pdf",
        ".ppt",
        ".pptx",
        ".png",
        ".jpg",
        ".jpeg",
        ".bmp",
        ".docx",
        ".doc",
        ".mp4",
        ".mov",
        ".avi",
        ".mkv",
        ".webm",
    }

    def __init__(self, upload_dir: Path, llm_client=None, audio_transcriber: OfflineAudioTranscriber | None = None):
        self.upload_dir = upload_dir
        self.llm_client = llm_client
        self.audio_transcriber = audio_transcriber or OfflineAudioTranscriber(workspace_dir=upload_dir)
        self.ocr_reader = None
        self._ocr_load_attempted = False
        self._logger = logging.getLogger(__name__)

    def save_and_parse(self, uploaded_file, session_id: str):
        raw_name = (uploaded_file.filename or "").strip()
        suffix = Path(raw_name).suffix.lower().strip()
        safe_stem = secure_filename(Path(raw_name).stem) or "upload"
        display_name = raw_name or f"{safe_stem}{suffix}"
        if suffix not in self.ALLOWED_SUFFIXES:
            raise ValueError(f"暂不支持的文件类型：{suffix or '无后缀'}")

        saved_name = f"{session_id}_{uuid.uuid4().hex}{suffix}"
        saved_path = self.upload_dir / saved_name
        uploaded_file.save(saved_path)
        parsed = self._parse_file(saved_path, display_name=display_name, enable_pdf_ocr=True)
        parsed["saved_name"] = saved_name
        return parsed

    def parse_existing_file(self, file_path: Path):
        path = Path(file_path)
        suffix = path.suffix.lower().strip()
        if suffix not in self.ALLOWED_SUFFIXES:
            raise ValueError(f"暂不支持的文件类型：{suffix or '无后缀'}")
        parsed = self._parse_file(path, display_name=path.name, enable_pdf_ocr=False)
        parsed["saved_name"] = path.name
        return parsed

    def _parse_file(self, path: Path, display_name: str, enable_pdf_ocr: bool):
        suffix = path.suffix.lower().strip()
        parser_map = {
            ".pdf": ("pdf", lambda p: self._extract_pdf(p, enable_ocr=enable_pdf_ocr)),
            ".pptx": ("ppt", self._extract_pptx),
            ".ppt": ("ppt", self._extract_legacy_ppt),
            ".docx": ("word", self._extract_docx),
            ".doc": ("word", self._extract_legacy_word),
            ".png": ("image", self._extract_image),
            ".jpg": ("image", self._extract_image),
            ".jpeg": ("image", self._extract_image),
            ".bmp": ("image", self._extract_image),
            ".mp4": ("video", self._extract_video),
            ".mov": ("video", self._extract_video),
            ".avi": ("video", self._extract_video),
            ".mkv": ("video", self._extract_video),
            ".webm": ("video", self._extract_video),
        }
        file_type, parser = parser_map[suffix]
        raw_text = parser(path)
        structured = self._build_structured_notes(original_name=display_name, file_type=file_type, text=raw_text)

        if not structured["clean_text"].strip():
            raise ValueError(f"{display_name} 未解析出可用内容，请确认文件未损坏，或尝试另存为可复制文本的 PDF / DOCX。")

        preview_lines = []
        if structured["summary"]:
            preview_lines.append(f"摘要：{structured['summary']}")
        if structured["key_points"]:
            preview_lines.append(f"要点：{'；'.join(structured['key_points'][:4])}")
        if structured["content_style"]:
            preview_lines.append(f"风格：{structured['content_style']}")
        preview_lines.append(structured["clean_text"][:400])

        return {
            "original_name": display_name,
            "file_type": file_type,
            "text": structured["clean_text"],
            "rag_text": structured["rag_text"],
            "summary": structured["summary"],
            "key_points": structured["key_points"],
            "knowledge_structure": structured["knowledge_structure"],
            "cases": structured["cases"],
            "content_style": structured["content_style"],
            "preview": "\n".join(preview_lines),
        }

    def _ensure_ocr_reader(self):
        if self.ocr_reader is not None:
            return self.ocr_reader
        if self._ocr_load_attempted:
            return None
        self._ocr_load_attempted = True
        try:
            self.ocr_reader = easyocr.Reader(["ch_sim", "en"], gpu=False)
        except Exception as error:
            self._logger.warning("EasyOCR 初始化失败，图像解析降级：%s", error)
            self.ocr_reader = None
        return self.ocr_reader

    def _extract_pdf(self, file_path: Path, enable_ocr: bool = True) -> str:
        text_parts = []
        try:
            reader = PdfReader(str(file_path))
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    pass
            for page in reader.pages:
                try:
                    page_text = (page.extract_text() or "").strip()
                    if page_text:
                        text_parts.append(page_text)
                except Exception:
                    continue
        except Exception as error:
            self._logger.warning("PyPDF2 解析失败：%s", error)

        if text_parts:
            return "\n".join(text_parts)

        try:
            import fitz
        except Exception:
            fitz = None

        if fitz is None:
            return ""

        try:
            document = fitz.open(str(file_path))
            fitz_text = []
            for page in document:
                page_text = (page.get_text("text") or "").strip()
                if page_text:
                    fitz_text.append(page_text)
            if fitz_text:
                return "\n".join(fitz_text)

            if not enable_ocr:
                return ""

            ocr_texts = []
            for page in document:
                pix = page.get_pixmap(dpi=220)
                image = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
                image = self._preprocess_image_for_ocr(image)
                text = self._run_ocr(image)
                if text:
                    ocr_texts.append(text)
            return "\n".join(ocr_texts)
        except Exception as error:
            self._logger.warning("PyMuPDF 解析失败：%s", error)
            return ""

    def _extract_docx(self, file_path: Path) -> str:
        try:
            document = DocxDocument(str(file_path))
        except Exception as error:
            self._logger.warning("DOCX 解析失败：%s", error)
            return ""

        blocks = []
        for paragraph in document.paragraphs:
            text = paragraph.text.strip()
            if text:
                blocks.append(text)
        for table in document.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    blocks.append(row_text)
        return "\n".join(blocks)

    def _extract_legacy_word(self, file_path: Path) -> str:
        return f"该文件为旧版 Word 文档：{file_path.name}。建议另存为 DOCX 后重新导入。"

    def _extract_pptx(self, file_path: Path) -> str:
        try:
            presentation = Presentation(str(file_path))
        except Exception as error:
            self._logger.warning("PPTX 解析失败：%s", error)
            return ""

        blocks = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_lines = []
            for shape in slide.shapes:
                text_chunks = []
                if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                    for paragraph in shape.text_frame.paragraphs:
                        text = "".join(run.text for run in paragraph.runs).strip() if paragraph.runs else paragraph.text.strip()
                        if text:
                            text_chunks.append(text)
                if getattr(shape, "has_table", False) and shape.has_table:
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            text_chunks.append(row_text)
                if text_chunks:
                    slide_lines.extend(text_chunks)

            notes_text = ""
            try:
                notes_frame = slide.notes_slide.notes_text_frame if slide.has_notes_slide else None
                notes_text = (notes_frame.text or "").strip() if notes_frame is not None else ""
            except Exception:
                notes_text = ""
            if notes_text:
                slide_lines.append(f"讲者备注：{notes_text}")

            if slide_lines:
                blocks.append(f"第{slide_index}页：" + "\n".join(slide_lines))

        return "\n\n".join(blocks)

    def _extract_legacy_ppt(self, file_path: Path) -> str:
        return f"该文件为旧版 PPT 文档：{file_path.name}。建议另存为 PPTX 后重新导入。"

    def _extract_image(self, file_path: Path) -> str:
        try:
            image = Image.open(file_path).convert("RGB")
        except Exception:
            return ""
        image = self._preprocess_image_for_ocr(image)
        return self._run_ocr(image)

    def _extract_video(self, file_path: Path) -> str:
        transcript = self.audio_transcriber.transcribe_video(file_path)
        if cv2 is None:
            return f"音频转写：{transcript}" if transcript else ""

        capture = cv2.VideoCapture(str(file_path))
        if not capture.isOpened():
            return f"音频转写：{transcript}" if transcript else ""

        frame_count = int(capture.get(cv2.CAP_PROP_FRAME_COUNT) or 0)
        sample_points = set()
        sample_total = 6
        if frame_count > 0:
            for index in range(sample_total):
                sample_points.add(int(index * max(frame_count - 1, 1) / max(sample_total - 1, 1)))
        else:
            sample_points = {0, 30, 60, 90, 120, 150}

        snippets = []
        frame_index = 0
        while True:
            ok, frame = capture.read()
            if not ok:
                break
            if frame_index in sample_points:
                image = Image.fromarray(cv2.cvtColor(frame, cv2.COLOR_BGR2RGB))
                image = self._preprocess_image_for_ocr(image)
                text = self._run_ocr(image)
                if text:
                    snippets.append(f"画面{frame_index}：{text}")
            frame_index += 1
        capture.release()

        parts = []
        if transcript:
            parts.append(f"音频转写：{transcript}")
        if snippets:
            parts.append("\n".join(snippets))
        return "\n".join(parts)

    def _preprocess_image_for_ocr(self, image: Image.Image):
        gray = image.convert("L")
        arr = np.array(gray)
        if arr.mean() < 120:
            arr = 255 - arr
        arr = np.where(arr > 170, 255, 0).astype("uint8")
        return Image.fromarray(arr)

    def _run_ocr(self, image: Image.Image):
        reader = self._ensure_ocr_reader()
        if reader is None:
            return ""
        try:
            result = reader.readtext(np.array(image), detail=0, paragraph=True)
            return "\n".join(str(item).strip() for item in result if str(item).strip())
        except Exception:
            return ""

    def _build_structured_notes(self, original_name: str, file_type: str, text: str):
        clean_text = self._normalize_text(text)
        knowledge_structure = self._extract_sentences(clean_text, ["定义", "概念", "原理", "步骤", "结构", "目标", "方法", "流程", "知识点", "公式", "定理"])
        cases = self._extract_sentences(clean_text, ["案例", "例如", "情境", "实验", "活动", "练习", "题", "应用", "任务", "例题"])
        content_style = self._infer_style(clean_text, file_type, original_name)
        key_points = []
        for item in knowledge_structure[:3] + cases[:3]:
            if item not in key_points:
                key_points.append(item)
        summary = self._build_summary(clean_text, key_points, content_style, original_name)
        compact_text = self._compact_text(clean_text, knowledge_structure, cases, content_style)
        rag_text = self._build_rag_text(clean_text, compact_text, knowledge_structure, cases)
        return {
            "clean_text": compact_text,
            "rag_text": rag_text,
            "summary": summary,
            "key_points": key_points[:6],
            "knowledge_structure": knowledge_structure[:6],
            "cases": cases[:6],
            "content_style": content_style,
        }

    def _normalize_text(self, text: str):
        text = str(text or "").replace("\x00", " ")
        text = re.sub(r"[ \t]+", " ", text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return text.strip()

    def _extract_sentences(self, text: str, keywords: list[str]):
        sentences = re.split(r"[。\n？！?!；;]", text)
        results = []
        for sentence in sentences:
            cleaned = sentence.strip()
            if len(cleaned) < 8:
                continue
            if any(keyword in cleaned for keyword in keywords):
                results.append(cleaned[:80])
        return results

    def _infer_style(self, text: str, file_type: str, original_name: str):
        sample = f"{original_name} {text[:1200]}".lower()
        if any(token in sample for token in ["实验", "活动单", "讨论", "互动", "任务"]):
            return "互动活动型，适合加入流程图、任务卡和分组页。"
        if any(token in sample for token in ["学术", "概念", "定义", "原理", "证明", "公式"]):
            return "学术讲解型，适合简洁排版、结构化标题和重点框。"
        if any(token in sample for token in ["案例", "场景", "应用", "项目", "真实问题", "例题"]):
            return "案例驱动型，适合情境导入、案例拆解和总结对照。"
        if file_type == "video":
            return "视听演示型，适合大图封面、步骤页、案例页，并可吸收讲解音轨内容。"
        return "通用教学型，适合目录-讲解-案例-小结结构。"

    def _build_summary(self, clean_text: str, key_points: list[str], content_style: str, original_name: str):
        if key_points:
            return f"{original_name} 的核心信息包括：{'；'.join(key_points[:3])}。建议采用{content_style}"
        if clean_text:
            return f"{original_name} 已提取到可用教学内容，可用于补充知识结构、案例和版式风格。"
        return f"{original_name} 未提取到足够文本。"

    def _compact_text(self, clean_text: str, knowledge_structure: list[str], cases: list[str], content_style: str):
        sections = []
        if knowledge_structure:
            sections.append("知识结构：" + "；".join(knowledge_structure[:4]))
        if cases:
            sections.append("案例素材：" + "；".join(cases[:4]))
        sections.append("排版风格：" + content_style)
        if not knowledge_structure and not cases and clean_text:
            sections.append("原始摘录：" + clean_text[:800])
        return "\n".join(sections)[:1600]

    def _build_rag_text(self, clean_text: str, compact_text: str, knowledge_structure: list[str], cases: list[str]):
        parts = [compact_text]
        if knowledge_structure:
            parts.append("知识点扩展：" + "；".join(knowledge_structure[:6]))
        if cases:
            parts.append("案例扩展：" + "；".join(cases[:6]))
        if clean_text:
            parts.append("原文片段：" + clean_text[:2200])
        return "\n".join(part for part in parts if part).strip()[:3600]
